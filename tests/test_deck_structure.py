"""簡報結構：封面、目錄、章節分隔頁、內容頁、結尾頁。

附件三的範例是 19 頁的完整簡報（封面 → 目錄 → 摘要 → 5 個章節 → 結尾），
不是一疊內容頁。這裡驗的是 renderer 真的組出這個骨架，以及**頁碼對得上**。

頁碼特別容易錯，也特別難發現：稽核 Excel 的工作表名是 `P.{頁碼}_{指標名}`
（FR-3.1），主管拿著 Excel 對照簡報。封面、目錄與每一張章節分隔頁都會讓
「第幾個內容頁」和「第幾張投影片」錯開，差一頁就是翻錯頁。
"""

import pytest
from pptx import Presentation

from ppt_generation.agents.narrative_writer import PageNarrative
from ppt_generation.agents.section_planner import (
    DEFAULT_CHAPTERS,
    SectionPlan,
    group_by_chapter,
)
from ppt_generation.core import config
from ppt_generation.data.metric_store import MetricSeries, MetricStore
from ppt_generation.output import renderer


TEMPLATE = config.TEMPLATE_PPTX


def _store():
    store = MetricStore(source_files=["fsc.xlsx"])
    store.add(
        MetricSeries(
            metric_key="流通卡數.value",
            name="流通卡數",
            categories=["中國信託", "國泰世華", "台新銀行"],
            series={"11412": [4_512_345.0, 4_102_887.0, 3_880_112.0]},
            unit="張",
        )
    )
    return store


def _sections(*pairs):
    return [
        SectionPlan(title=title, intent="", chapter=chapter)
        for chapter, title in pairs
    ]


def _bundles(sections):
    return [
        renderer.PageBundle(
            section,
            None,
            PageNarrative(
                slide_title=section.title,
                headline="標題",
                bullets=["重點"],
                page_number=section.page_number,
            ),
        )
        for section in sections
    ]


# ---------------------------------------------------------------------------
# 章節收攏
# ---------------------------------------------------------------------------
def test_pages_of_the_same_chapter_are_grouped():
    """LLM 可能把同章節的頁面排散，章節分隔頁就會出現兩次。"""
    sections = _sections(
        ("市場整體概況", "規模"),
        ("風險與警訊", "逾期率"),
        ("市場整體概況", "成長"),
    )

    grouped = group_by_chapter(sections)

    assert [section.title for section in grouped] == ["規模", "成長", "逾期率"]


def test_chapter_order_follows_first_appearance():
    sections = _sections(
        ("風險與警訊", "逾期率"),
        ("市場整體概況", "規模"),
    )

    grouped = group_by_chapter(sections)

    assert [section.chapter for section in grouped] == [
        "風險與警訊",
        "市場整體概況",
    ]


def test_default_chapters_match_the_spec():
    """FR-2.6 明列的八章節。改動這份清單等於改規格，要有意識。"""
    assert len(DEFAULT_CHAPTERS) == 8
    assert DEFAULT_CHAPTERS[0] == "Executive Summary"
    assert "未來趨勢推測" in DEFAULT_CHAPTERS


# ---------------------------------------------------------------------------
# 頁碼指派
# ---------------------------------------------------------------------------
def test_page_numbers_account_for_front_matter_and_dividers():
    sections = _sections(
        ("市場整體概況", "規模"),
        ("市場整體概況", "成長"),
        ("風險與警訊", "逾期率"),
    )

    renderer.assign_page_numbers(sections)

    # 1 封面、2 目錄、3 章節頁、4 規模、5 成長、6 章節頁、7 逾期率
    assert [section.page_number for section in sections] == [4, 5, 7]


def test_page_numbers_skip_agenda_when_disabled():
    sections = _sections(("市場整體概況", "規模"))

    renderer.assign_page_numbers(sections, include_agenda=False)

    assert sections[0].page_number == 3


def test_page_numbers_without_chapters_have_no_agenda_or_divider():
    """沒有章節就沒有目錄可列，render_deck 也不會產目錄頁，兩邊必須一致。"""
    sections = [SectionPlan(title="單頁", intent="")]

    renderer.assign_page_numbers(sections)

    assert sections[0].page_number == 2


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_page_numbers_match_actual_slide_positions(tmp_path):
    """
    這是頁碼指派唯一真正重要的斷言：算出來的頁碼要等於實際投影片序號。

    兩邊各自實作（一邊推算、一邊組裝），所以必須交叉驗證，否則其中一邊
    改了版面順序卻沒改另一邊，稽核 Excel 就會整份對錯頁。
    """
    sections = _sections(
        ("市場整體概況", "規模"),
        ("市場整體概況", "成長"),
        ("風險與警訊", "逾期率"),
    )
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    renderer.render_deck(
        _bundles(sections), _store(), output_path=output
    )

    presentation = Presentation(str(output))
    titles = [
        slide.shapes.title.text_frame.text if slide.shapes.title else ""
        for slide in presentation.slides
    ]

    for section in sections:
        assert titles[section.page_number - 1] == section.title


# ---------------------------------------------------------------------------
# 版面組裝
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_deck_has_cover_agenda_dividers_and_closing(tmp_path):
    sections = _sections(
        ("市場整體概況", "規模"),
        ("風險與警訊", "逾期率"),
    )
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    report = renderer.render_deck(
        _bundles(sections),
        _store(),
        output_path=output,
        deck_title="信用卡市場分析",
    )

    # 封面 1 + 目錄 1 + 章節頁 2 + 內容頁 2 + 結論 1 + 結尾 1
    assert report.slide_count == 8
    assert report.page_count == 2
    assert report.divider_count == 2
    assert report.conclusion_page is True
    assert report.chapters == ["市場整體概況", "風險與警訊"]


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_cover_title_is_overwritten(tmp_path):
    sections = _sections(("市場整體概況", "規模"))
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    renderer.render_deck(
        _bundles(sections),
        _store(),
        output_path=output,
        deck_title="信用卡市場分析與經營洞察",
    )

    presentation = Presentation(str(output))

    assert (
        presentation.slides[0].shapes.title.text_frame.text
        == "信用卡市場分析與經營洞察"
    )


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_agenda_lists_every_chapter_once(tmp_path):
    """目錄與章節分隔頁不一致是最傷信任的瑕疵，所以目錄只由章節清單產生。"""
    sections = _sections(
        ("市場整體概況", "規模"),
        ("市場整體概況", "成長"),
        ("風險與警訊", "逾期率"),
    )
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    renderer.render_deck(_bundles(sections), _store(), output_path=output)

    presentation = Presentation(str(output))
    agenda_text = "\n".join(
        shape.text_frame.text
        for shape in presentation.slides[1].shapes
        if shape.has_text_frame
    )

    assert agenda_text.count("市場整體概況") == 1
    assert agenda_text.count("風險與警訊") == 1


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_dividers_use_the_template_section_layout(tmp_path):
    sections = _sections(("市場整體概況", "規模"))
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    renderer.render_deck(_bundles(sections), _store(), output_path=output)

    presentation = Presentation(str(output))
    divider = presentation.slides[2]

    assert divider.slide_layout.name == renderer.SECTION_LAYOUT_NAME
    # 章節頁除了章節名，還帶一行 CHAPTER 編號（對齊附件三）。
    assert divider.shapes.title.text_frame.text == "CHAPTER 01\n市場整體概況"


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_closing_page_uses_the_template_closing_layout(tmp_path):
    sections = _sections(("市場整體概況", "規模"))
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    renderer.render_deck(_bundles(sections), _store(), output_path=output)

    presentation = Presentation(str(output))
    closing = presentation.slides[-1]

    assert closing.slide_layout.name == renderer.CLOSING_LAYOUT_NAME


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_agenda_and_closing_can_be_turned_off(tmp_path):
    sections = _sections(("市場整體概況", "規模"))
    renderer.assign_page_numbers(sections, include_agenda=False)

    output = tmp_path / "deck.pptx"
    report = renderer.render_deck(
        _bundles(sections),
        _store(),
        output_path=output,
        include_agenda=False,
        include_conclusion=False,
        include_closing=False,
    )

    # 封面 1 + 章節頁 1 + 內容頁 1
    assert report.slide_count == 3
    assert report.conclusion_page is False


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_pages_without_chapter_get_no_divider(tmp_path):
    sections = [SectionPlan(title="單頁", intent="")]
    renderer.assign_page_numbers(sections)

    output = tmp_path / "deck.pptx"
    report = renderer.render_deck(
        _bundles(sections), _store(), output_path=output
    )

    assert report.divider_count == 0
    assert report.chapters == []
    # 封面 1 + 內容頁 1 + 結尾 1（無章節即無目錄，也無結論頁可收）
    assert report.slide_count == 3
    assert report.conclusion_page is False

"""雙軸圖（長條 + 折線掛次軸）。

附件三 P.5/P.6 是命題方的主打頁：流通卡數（張）與當月簽帳金額（百萬元）
量級差一個數量級，非用次軸不可。python-pptx 沒有雙軸的高階 API，
規格書把它列為 M0 唯一出場條件。

風險不在「畫得出來」，在**畫出來之後右鍵編輯資料還是對的**。自己組 chart XML
手填 <c:numCache> 會讓快取有值、內嵌 workbook 空白——那就是附件三的病灶本身。
所以這裡驗三件事：折線真的掛在次軸、內嵌 workbook 兩個系列都完整、
兩份副本逐格相同。
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_generation.charts.chart_builder import (
    CHART_SKILLS,
    ComboSpec,
    add_combo_chart,
)
from ppt_generation.verification import verify_chart_consistency as vcc


CATEGORIES = ["11401", "11402", "11403", "11404"]
CARDS = [58_723_456.0, 58_901_233.0, 59_120_876.0, 59_455_002.0]
AMOUNT = [312_456.7, 298_331.2, 341_902.5, 355_128.9]


def _slide():
    presentation = Presentation()
    layout = presentation.slide_layouts[6]  # 空白版面
    return presentation, presentation.slides.add_slide(layout)


def _spec():
    return ComboSpec(
        title="市場規模趨勢 — 流通卡數與簽帳金額",
        categories=CATEGORIES,
        series={"流通卡數": CARDS, "當月簽帳金額": AMOUNT},
        line_series_names=("當月簽帳金額",),
    )


def _build():
    presentation, slide = _slide()
    chart = add_combo_chart(
        slide,
        _spec(),
        left=Inches(0.5),
        top=Inches(1.0),
        width=Inches(9.0),
        height=Inches(5.0),
    )
    return presentation, chart


def _plot_area(chart):
    return chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))


# ---------------------------------------------------------------------------
# 結構
# ---------------------------------------------------------------------------
def test_registered_as_a_skill():
    """圖表 Agent 只透過 skill 名稱呼叫，沒註冊等於做了也用不到。"""
    assert CHART_SKILLS["combo"] is add_combo_chart


def test_both_plot_types_exist():
    _, chart = _build()
    plot_area = _plot_area(chart)

    assert plot_area.find(qn("c:barChart")) is not None
    assert plot_area.find(qn("c:lineChart")) is not None


def test_series_are_split_between_the_two_plots():
    _, chart = _build()
    plot_area = _plot_area(chart)

    bar_series = plot_area.find(qn("c:barChart")).findall(qn("c:ser"))
    line_series = plot_area.find(qn("c:lineChart")).findall(qn("c:ser"))

    assert len(bar_series) == 1
    assert len(line_series) == 1


def test_line_chart_uses_a_separate_axis_pair():
    """次軸的 axId 必須與主軸不同，否則折線會被壓回主軸刻度。"""
    _, chart = _build()
    plot_area = _plot_area(chart)

    bar_ids = {
        element.get("val")
        for element in plot_area.find(qn("c:barChart")).findall(qn("c:axId"))
    }
    line_ids = {
        element.get("val")
        for element in plot_area.find(qn("c:lineChart")).findall(qn("c:axId"))
    }

    assert len(bar_ids) == 2
    assert len(line_ids) == 2
    assert bar_ids.isdisjoint(line_ids)


def test_every_axis_id_is_declared_as_an_axis():
    """圖形群組引用了不存在的 axId，PowerPoint 會要求修復檔案。"""
    _, chart = _build()
    plot_area = _plot_area(chart)

    declared = {
        axis.find(qn("c:axId")).get("val")
        for tag in ("c:catAx", "c:valAx")
        for axis in plot_area.findall(qn(tag))
    }
    referenced = {
        element.get("val")
        for tag in ("c:barChart", "c:lineChart")
        for element in plot_area.find(qn(tag)).findall(qn("c:axId"))
    }

    assert referenced <= declared
    assert len(declared) == 4


def test_secondary_value_axis_is_on_the_right():
    _, chart = _build()
    plot_area = _plot_area(chart)

    positions = [
        axis.find(qn("c:axPos")).get("val")
        for axis in plot_area.findall(qn("c:valAx"))
    ]

    assert sorted(positions) == ["l", "r"]


def test_secondary_category_axis_is_hidden():
    """兩組類別軸都顯示的話，圖底下會出現兩排一樣的月份標籤。"""
    _, chart = _build()
    plot_area = _plot_area(chart)

    deleted = [
        axis.find(qn("c:delete")).get("val")
        for axis in plot_area.findall(qn("c:catAx"))
    ]

    assert sorted(deleted) == ["0", "1"]


def test_line_series_has_no_bar_only_children():
    """invertIfNegative 是長條專屬元素，留在折線 ser 裡會違反 schema。"""
    _, chart = _build()
    line_ser = (
        _plot_area(chart).find(qn("c:lineChart")).find(qn("c:ser"))
    )

    assert line_ser.find(qn("c:invertIfNegative")) is None


def test_line_series_keeps_its_values():
    _, chart = _build()
    line_ser = (
        _plot_area(chart).find(qn("c:lineChart")).find(qn("c:ser"))
    )

    assert line_ser.find(qn("c:val")) is not None
    assert line_ser.find(qn("c:cat")) is not None


# ---------------------------------------------------------------------------
# 資料一致性（這才是重點）
# ---------------------------------------------------------------------------
def test_embedded_workbook_keeps_both_series():
    """右鍵「編輯資料」看到的表必須含兩個系列，不能只剩長條那一個。"""
    _, chart = _build()

    embedded = vcc.read_embedded_series(chart)

    assert set(embedded) == {"流通卡數", "當月簽帳金額"}


def test_chart_cache_matches_embedded_workbook(tmp_path):
    """畫面顯示值 ↔ 內嵌 workbook 逐格相同（T1 的①②兩方）。"""
    presentation, _ = _build()
    path = tmp_path / "combo.pptx"
    presentation.save(str(path))

    report = vcc.verify(path)

    assert report.passed, [c.describe_failure() for c in report.failures]


def test_verification_covers_the_secondary_axis_series(tmp_path):
    """驗證只看 plots[0] 的話，次軸系列會靜默漏驗——比對數要是 2 不是 1。"""
    presentation, _ = _build()
    path = tmp_path / "combo.pptx"
    presentation.save(str(path))

    report = vcc.verify(path)
    names = {comparison.series_name for comparison in report.comparisons}

    assert names == {"流通卡數", "當月簽帳金額"}


def test_values_survive_the_plot_move(tmp_path):
    """搬動 c:ser 只改「用什麼圖形畫」，數值節點不得被動到。"""
    presentation, _ = _build()
    path = tmp_path / "combo.pptx"
    presentation.save(str(path))

    report = vcc.verify(path)
    by_name = {
        comparison.series_name: comparison.chart_cache
        for comparison in report.comparisons
    }

    assert by_name["流通卡數"] == pytest.approx(CARDS)
    assert by_name["當月簽帳金額"] == pytest.approx(AMOUNT)


# ---------------------------------------------------------------------------
# 防呆
# ---------------------------------------------------------------------------
def test_rejects_missing_line_series_names():
    _, slide = _slide()
    spec = ComboSpec(
        title="缺折線系列",
        categories=CATEGORIES,
        series={"流通卡數": CARDS},
    )

    with pytest.raises(ValueError):
        add_combo_chart(slide, spec)


def test_rejects_unknown_line_series_name():
    _, slide = _slide()
    spec = ComboSpec(
        title="系列名打錯",
        categories=CATEGORIES,
        series={"流通卡數": CARDS, "當月簽帳金額": AMOUNT},
        line_series_names=("簽帳金額",),
    )

    with pytest.raises(ValueError):
        add_combo_chart(slide, spec)


def test_rejects_when_all_series_become_lines():
    """全部變折線就沒有主軸長條，這種圖應該一開始就用 line。"""
    _, slide = _slide()
    spec = ComboSpec(
        title="全是折線",
        categories=CATEGORIES,
        series={"流通卡數": CARDS, "當月簽帳金額": AMOUNT},
        line_series_names=("流通卡數", "當月簽帳金額"),
    )

    with pytest.raises(ValueError):
        add_combo_chart(slide, spec)


def test_three_series_split_two_to_one():
    """兩個以上的折線系列共用同一組次軸，不會各自再長一組軸。"""
    _, slide = _slide()
    spec = ComboSpec(
        title="三系列",
        categories=CATEGORIES,
        series={
            "流通卡數": CARDS,
            "當月簽帳金額": AMOUNT,
            "有效卡數": [value * 0.8 for value in CARDS],
        },
        line_series_names=("當月簽帳金額", "有效卡數"),
    )

    chart = add_combo_chart(slide, spec)
    plot_area = _plot_area(chart)

    assert len(plot_area.find(qn("c:barChart")).findall(qn("c:ser"))) == 1
    assert len(plot_area.find(qn("c:lineChart")).findall(qn("c:ser"))) == 2
    assert len(plot_area.findall(qn("c:valAx"))) == 2

"""簡報風格：附件三的版面語彙與「重點用黃色」的落地。

這裡驗三件在簡報上肉眼可見、但很容易在重構中默默壞掉的事：

1. **黃色只標在值上。** 由 MetricStore 代入的數字要有 `a:highlight`，
   敘事文字不能有。整段塗黃就不是重點而是塗滿，而漏掉則失去「哪些字是
   系統算的」這條線索。
2. **重點訊息帶存在且只有一條。** 附件三每頁都有一句結論句在標題下方，
   它是「一頁一結論」的載體。
3. **結論頁把每個章節收成一行。** 章節規劃有結論，是簡報而不是資料集。

另外驗頁尾註記不可張冠李戴：原生表格沒有內嵌工作簿，右鍵不會有
「編輯資料」，沿用圖表那句就是在簡報上寫下不成立的承諾。
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from ppt_generation.agents.narrative_writer import (
    MIN_BULLETS,
    PageNarrative,
    check_narrative,
)
from ppt_generation.agents.section_planner import CONCLUSION_CHAPTER, SectionPlan
from ppt_generation.charts.chart_builder import ChartSpec
from ppt_generation.charts.chart_planner import ChartPlan, ResolvedChart
from ppt_generation.charts.table_builder import TableSpec
from ppt_generation.core import config, placeholders
from ppt_generation.data.metric_store import MetricSeries, MetricStore
from ppt_generation.output import renderer, theme


TEMPLATE = config.TEMPLATE_PPTX


def _store():
    store = MetricStore(source_files=["fsc.xlsx"])
    store.add(
        MetricSeries(
            metric_key="流通卡數.value",
            name="流通卡數",
            categories=["中國信託", "國泰世華", "台新銀行"],
            series={"11412": [4_512_345.0, 4_102_887.0, 3_880_112.0]},
            unit="張",
        )
    )
    return store


def _narrative(page_number=4, title="規模"):
    return PageNarrative(
        page_number=page_number,
        slide_title=title,
        headline="龍頭 {{流通卡數.value|11412|max_category}} 規模領先，差距短期難以收斂",
        bullets=[
            "領先者達 {{流通卡數.value|11412|max}}，明顯拉開與其餘業者的距離，"
            "反映規模效應仍在累積",
            "末位者僅 {{流通卡數.value|11412|min}}，缺乏可持續的規模優勢，"
            "須從差異化定位切入",
            "全體平均 {{流通卡數.value|11412|avg}}，位於平均之下者"
            "需重新配置行銷資源與通路組合，避免無效投放",
        ],
    )


def _chart(skill_name="column", page_number=4):
    spec = ChartSpec(
        title="流通卡數",
        categories=["中國信託", "國泰世華", "台新銀行"],
        series={"11412": [4_512_345.0, 4_102_887.0, 3_880_112.0]},
    )

    if skill_name == "table":
        spec = TableSpec(
            title="流通卡數",
            categories=spec.categories,
            series=spec.series,
        )

    plan = ChartPlan(
        slide_title="規模",
        chart_title="流通卡數",
        metric_key="流通卡數.value",
        chart_type=skill_name,
        page_number=page_number,
    )

    return ResolvedChart(
        plan=plan,
        metric=_store().get("流通卡數.value"),
        spec=spec,
        skill_name=skill_name,
    )


def _render(tmp_path, *, skill_name="column", chapters=("市場整體概況",)):
    sections = [
        SectionPlan(title="規模", intent="", chapter=chapter)
        for chapter in chapters
    ]
    renderer.assign_page_numbers(sections)

    bundles = [
        renderer.PageBundle(
            section,
            _chart(skill_name, section.page_number),
            _narrative(section.page_number, section.title),
        )
        for section in sections
    ]

    output = tmp_path / "deck.pptx"
    report = renderer.render_deck(bundles, _store(), output_path=output)

    return Presentation(str(output)), report


# ---------------------------------------------------------------------------
# 佔位符片段
# ---------------------------------------------------------------------------
def test_render_segments_marks_only_the_looked_up_values():
    segments, errors = placeholders.render_segments(
        "龍頭市占 {{流通卡數.value|11412|max}} 居首", _store()
    )

    assert errors == []
    assert [segment.from_metric for segment in segments] == [False, True, False]
    assert segments[1].metric_key == "流通卡數.value"


def test_render_segments_concatenate_to_render_text():
    """兩個入口共用同一組查表函式，串起來必須逐字相同，否則就有兩種數字。"""
    text = "{{流通卡數.value|11412|max}} 對 {{流通卡數.value|11412|min}}"
    segments, _ = placeholders.render_segments(text, _store())
    rendered, _ = placeholders.render_text(text, _store())

    assert "".join(segment.text for segment in segments) == rendered


def test_render_segments_keeps_failed_placeholder_visible():
    segments, errors = placeholders.render_segments(
        "市占 {{不存在的指標|系列|max}}", _store(), strict=False
    )

    assert errors
    assert "{{不存在的指標|系列|max}}" in "".join(s.text for s in segments)
    assert not any(segment.from_metric for segment in segments)


# ---------------------------------------------------------------------------
# 黃色標示
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_only_metric_values_are_highlighted_yellow(tmp_path):
    presentation, _ = _render(tmp_path)
    content = presentation.slides[3]

    highlighted: list[str] = []
    plain: list[str] = []

    for shape in content.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                element = run._r.find(qn("a:rPr"))
                marker = (
                    element.find(qn("a:highlight"))
                    if element is not None
                    else None
                )

                if marker is None:
                    plain.append(run.text)
                    continue

                srgb = marker.find(qn("a:srgbClr"))
                assert srgb.get("val") == f"{theme.HIGHLIGHT}"
                highlighted.append(run.text)

    assert highlighted, "代入的數值完全沒有被標示"
    # 被標黃的每一段都必須是查表得來的值，不能是敘事文字。
    for text in highlighted:
        assert any(char.isdigit() for char in text) or text == "中國信託"

    assert any("領先者" in text for text in plain)


# ---------------------------------------------------------------------------
# 重點訊息帶
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_key_message_bar_carries_the_headline_once(tmp_path):
    presentation, _ = _render(tmp_path)
    content = presentation.slides[3]

    bars = [
        shape
        for shape in content.shapes
        if shape.has_text_frame
        and shape.text_frame.text.startswith(theme.KEY_MESSAGE_PREFIX)
    ]

    assert len(bars) == 1
    bar = bars[0]

    assert "規模領先" in bar.text_frame.text
    assert bar.fill.fore_color.rgb == theme.KEY_BAR_FILL
    # headline 只在訊息帶出現，不重複進正文，否則同一句話一頁講兩次。
    body = next(
        shape
        for shape in content.shapes
        if shape.has_text_frame
        and shape.text_frame.text.startswith(theme.ACTION_BULLET_PREFIX)
    )
    assert "規模領先" not in body.text_frame.text


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_key_message_bar_does_not_overlap_the_chart(tmp_path):
    presentation, _ = _render(tmp_path)
    content = presentation.slides[3]

    bar = next(
        shape
        for shape in content.shapes
        if shape.has_text_frame
        and shape.text_frame.text.startswith(theme.KEY_MESSAGE_PREFIX)
    )
    chart = next(shape for shape in content.shapes if shape.has_chart)

    assert bar.top + bar.height <= chart.top


# ---------------------------------------------------------------------------
# 頁尾註記
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_chart_pages_state_the_right_click_check(tmp_path):
    presentation, _ = _render(tmp_path, skill_name="column")
    texts = [
        shape.text_frame.text
        for shape in presentation.slides[3].shapes
        if shape.has_text_frame
    ]

    assert any(theme.CHART_FOOTNOTE == text for text in texts)


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_table_pages_do_not_promise_right_click_editing(tmp_path):
    """原生表格沒有內嵌工作簿，右鍵不會有「編輯資料」，不能沿用圖表那句。"""
    presentation, _ = _render(tmp_path, skill_name="table")
    texts = [
        shape.text_frame.text
        for shape in presentation.slides[3].shapes
        if shape.has_text_frame
    ]

    assert any(theme.TABLE_FOOTNOTE == text for text in texts)
    assert all(theme.CHART_FOOTNOTE != text for text in texts)


# ---------------------------------------------------------------------------
# 結論頁
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_conclusion_page_collects_one_line_per_chapter(tmp_path):
    presentation, report = _render(
        tmp_path, chapters=("市場整體概況", "風險與警訊")
    )

    assert report.conclusion_page is True

    # 倒數第二張是結論頁（最後一張是結尾頁）。
    conclusion = presentation.slides[-2]

    assert conclusion.shapes.title.text_frame.text == CONCLUSION_CHAPTER

    body = "\n".join(
        shape.text_frame.text
        for shape in conclusion.shapes
        if shape.has_text_frame and not shape == conclusion.shapes.title
    )

    for chapter in ("市場整體概況", "風險與警訊"):
        assert body.count(chapter) == 1

    # 每一行都指回它出自哪一頁，主管問「這句從哪來」答案就在簡報裡。
    assert "（P." in body


@pytest.mark.skipif(not TEMPLATE.exists(), reason="找不到 source/template.pptx")
def test_agenda_ends_with_the_conclusion_chapter(tmp_path):
    presentation, _ = _render(tmp_path, chapters=("市場整體概況",))
    agenda = "\n".join(
        shape.text_frame.text
        for shape in presentation.slides[1].shapes
        if shape.has_text_frame
    )

    assert CONCLUSION_CHAPTER in agenda


# ---------------------------------------------------------------------------
# 敘事份量
# ---------------------------------------------------------------------------
def test_narrative_with_enough_substance_passes():
    issues = check_narrative(
        _narrative(), _store(), {"流通卡數.value"}
    )

    assert issues == []


def test_too_few_bullets_is_rejected():
    narrative = _narrative()
    narrative.bullets = narrative.bullets[:1]

    issues = check_narrative(narrative, _store(), {"流通卡數.value"})

    assert any(f"下限 {MIN_BULLETS}" in issue for issue in issues)


def test_stub_bullets_are_rejected():
    """「本頁無可引用指標」這種交差式回應要在敘事階段就被擋下。"""
    narrative = _narrative()
    narrative.bullets = ["資料觀察", "數值偏高", "值得注意"]

    issues = check_narrative(narrative, _store(), {"流通卡數.value"})

    assert any("低於下限" in issue for issue in issues)


def test_narrative_without_any_placeholder_is_rejected():
    narrative = PageNarrative(
        page_number=4,
        slide_title="規模",
        headline="市場呈現高度集中的競爭格局，短期難以改變",
        bullets=[
            "領先業者掌握主要份額，後段業者缺乏規模優勢與差異化定位可言",
            "中段業者若要突圍，需從單卡消費力而非發卡量切入市場競爭",
            "整體而言，資源配置應與市場競爭強度相稱，避免無效投放預算",
        ],
    )

    issues = check_narrative(narrative, _store(), {"流通卡數.value"})

    assert any("沒有引用任何指標佔位符" in issue for issue in issues)

"""原生表格與熱力圖（FR-2.4 / FR-2.3）。

附件三的問題之一是「表格以幾何圖形和文字方塊堆疊，非真正表格」。所以這裡
第一件要驗的就是產出物真的是 PPT table 物件（``shape.has_table``），而不是
看起來像表格的一堆方塊。

第二件是數字。表格沒有內嵌 workbook，右鍵沒有「編輯資料」——這是 PowerPoint
表格的本性，不是缺陷。但它的後果是表格數字少一份可互相印證的副本，因此
格式化必須只有一個入口（``format_value``），且要能被反向解析回來與稽核
Excel 比對。這條路徑斷了，表格就變成系統裡唯一沒人看守的數字。
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_generation.charts import table_builder
from ppt_generation.charts.table_builder import (
    HEATMAP_HIGH,
    HEATMAP_LOW,
    MAX_TABLE_ROWS,
    TABLE_SKILLS,
    TableBuildError,
    TableSpec,
    add_heatmap_table,
    add_native_table,
    format_value,
    heatmap_color,
    parse_value,
)


BANKS = ["中國信託", "國泰世華", "台新銀行", "總計"]
CARDS = [4_512_345.0, 4_102_887.0, 3_880_112.0, 58_723_456.0]
SHARES = [15.6369, 14.2131, 13.4459, None]


def _slide():
    presentation = Presentation()
    return presentation, presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )


def _spec(**overrides):
    defaults = dict(
        title="流通卡數與市占率",
        categories=BANKS,
        series={"流通卡數": CARDS},
        unit="張",
        emphasize_rows=("總計",),
    )
    defaults.update(overrides)
    return TableSpec(**defaults)


# ---------------------------------------------------------------------------
# 是不是真的原生表格
# ---------------------------------------------------------------------------
def test_registered_as_skills():
    assert TABLE_SKILLS["table"] is add_native_table
    assert TABLE_SKILLS["heatmap"] is add_heatmap_table


def test_produces_a_real_pptx_table():
    """附件三的病灶：用文字方塊拼出表格的樣子。這裡必須是 table 物件。"""
    _, slide = _slide()

    add_native_table(slide, _spec())

    tables = [shape for shape in slide.shapes if shape.has_table]
    assert len(tables) == 1


def test_shape_dimensions_match_the_spec():
    _, slide = _slide()

    table = add_native_table(slide, _spec())

    # 多一列表頭、多一欄列標籤
    assert len(table.rows) == len(BANKS) + 1
    assert len(table.columns) == 2


def test_header_row_holds_row_header_and_series_names():
    _, slide = _slide()

    table = add_native_table(
        slide,
        _spec(series={"流通卡數": CARDS, "市占率": [10.0, 9.0, 8.0, None]}),
    )
    header = [cell.text for cell in table.rows[0].cells]

    assert header == ["項目", "流通卡數", "市占率"]


def test_row_labels_are_written_in_order():
    _, slide = _slide()

    table = add_native_table(slide, _spec())
    labels = [row.cells[0].text for row in list(table.rows)[1:]]

    assert labels == BANKS


# ---------------------------------------------------------------------------
# 數值格式化（唯一入口）
# ---------------------------------------------------------------------------
@pytest.mark.parametrize(
    "value,unit,expected",
    [
        (4_512_345.0, "張", "4,512,345"),
        (15.6369, "%", "15.64%"),
        (3.0, "名", "3"),
        (1234.5, None, "1,234.50"),
        (None, "張", "—"),
        (None, "%", "—"),
    ],
)
def test_format_value(value, unit, expected):
    assert format_value(value, unit) == expected


def test_missing_value_is_a_dash_not_zero():
    """缺值印成 0 就是在圖上宣告一個不存在的觀測值（附件三錯誤的一種形式）。"""
    assert format_value(None) == "—"
    assert format_value(0.0) == "0"


@pytest.mark.parametrize(
    "text,expected",
    [
        ("4,512,345", 4_512_345.0),
        ("15.64%", 15.64),
        ("—", None),
        ("", None),
        ("不是數字", None),
    ],
)
def test_parse_value_round_trips(text, expected):
    assert parse_value(text) == expected


def test_format_and_parse_are_a_pair():
    """驗證要靠這一對函式互為反向；任一邊改了另一邊沒跟上，比對就會誤判。"""
    for value, unit in ((4_512_345.0, "張"), (15.64, "%"), (3.0, "名")):
        assert parse_value(format_value(value, unit), unit) == value


def test_missing_values_render_as_dash_in_cells():
    _, slide = _slide()

    table = add_native_table(
        slide,
        _spec(series={"市占率": SHARES}, unit="%"),
    )
    last_row = list(table.rows)[-1]

    assert last_row.cells[1].text == "—"


# ---------------------------------------------------------------------------
# 熱力圖色階
# ---------------------------------------------------------------------------
def test_heatmap_color_endpoints():
    assert heatmap_color(0.0, 0.0, 10.0) == HEATMAP_LOW
    assert heatmap_color(10.0, 0.0, 10.0) == HEATMAP_HIGH


def test_heatmap_color_is_monotonic():
    colors = [
        heatmap_color(value, 0.0, 10.0) for value in (1.0, 5.0, 9.0)
    ]
    # 高值往深色走，所以 R 通道遞減
    assert colors[0][0] > colors[1][0] > colors[2][0]


def test_heatmap_skips_missing_values():
    """缺值不上色。塗成最淺色等於宣稱它是全表最小值。"""
    assert heatmap_color(None, 0.0, 10.0) is None


def test_heatmap_flat_data_uses_low_color():
    """全表同值時色階沒有資訊量，用深色會讓讀者以為每格都是高點。"""
    assert heatmap_color(5.0, 5.0, 5.0) == HEATMAP_LOW


def test_heatmap_fills_cells():
    _, slide = _slide()

    table = add_heatmap_table(
        slide,
        _spec(series={"流通卡數": CARDS}),
    )
    data_cell = list(table.rows)[1].cells[1]

    assert data_cell.fill.fore_color.rgb is not None


def test_plain_table_does_not_fill_data_cells():
    _, slide = _slide()

    table = add_native_table(slide, _spec())
    data_cell = list(table.rows)[1].cells[1]

    # 未上色的儲存格 fill 型別不是 solid
    assert data_cell.fill.type != 1


def test_heatmap_does_not_change_any_value():
    """色階只是視覺編碼。上色後數字若變了，就是把呈現混進了計算。"""
    _, slide = _slide()

    plain = add_native_table(slide, _spec())
    _, slide2 = _slide()
    heat = add_heatmap_table(slide2, _spec())

    def texts(table):
        return [
            row.cells[1].text for row in list(table.rows)[1:]
        ]

    assert texts(plain) == texts(heat)


# ---------------------------------------------------------------------------
# 防呆
# ---------------------------------------------------------------------------
def test_rejects_empty_rows():
    _, slide = _slide()

    with pytest.raises(TableBuildError):
        add_native_table(slide, _spec(categories=[], series={"流通卡數": []}))


def test_rejects_empty_columns():
    _, slide = _slide()

    with pytest.raises(TableBuildError):
        add_native_table(slide, _spec(series={}))


def test_rejects_too_many_rows():
    """一頁塞 33 家銀行的表沒有人讀得完，應該先取 Top N 或改用圖表。"""
    _, slide = _slide()
    count = MAX_TABLE_ROWS + 1

    with pytest.raises(TableBuildError):
        add_native_table(
            slide,
            _spec(
                categories=[f"銀行{index}" for index in range(count)],
                series={"流通卡數": [1.0] * count},
            ),
        )


def test_heatmap_rejects_all_missing_values():
    _, slide = _slide()

    with pytest.raises(TableBuildError):
        add_heatmap_table(
            slide,
            _spec(series={"流通卡數": [None] * len(BANKS)}),
        )


# ---------------------------------------------------------------------------
# 與驗證模組的接點
# ---------------------------------------------------------------------------
def test_table_can_be_read_back_for_verification(tmp_path):
    """驗證模組要能把表格讀回數值，否則表格上的數字沒人看守。"""
    from ppt_generation.verification import verify_chart_consistency as vcc

    presentation, slide = _slide()
    add_native_table(slide, _spec())

    path = tmp_path / "table.pptx"
    presentation.save(str(path))

    reloaded = Presentation(str(path))
    shape = next(
        shape for shape in reloaded.slides[0].shapes if shape.has_table
    )
    columns, labels, values = vcc.read_table_shape(shape)

    assert columns == ["流通卡數"]
    assert labels == BANKS
    assert values["流通卡數"] == CARDS


def test_table_without_audit_excel_is_not_considered_passed(tmp_path):
    """表格沒有第二份副本，缺稽核 Excel 就是「無法驗證」而非「通過」。"""
    from ppt_generation.verification import verify_chart_consistency as vcc

    presentation, slide = _slide()
    add_native_table(slide, _spec())

    path = tmp_path / "table.pptx"
    presentation.save(str(path))

    report = vcc.verify(path)

    assert report.comparisons
    assert not report.passed


def test_module_exposes_a_single_formatter():
    """表格文字只能有一個產生入口，否則系統裡就有第二個數字來源。"""
    assert table_builder.format_value is format_value

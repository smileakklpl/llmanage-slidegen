"""Cross-domain acceptance through backend → core → ppt_generation."""

from __future__ import annotations

import hashlib
import json
import sys
from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation

from core.contracts.generation import GenerationRequest
from core.generation_orchestrator import (
    GenerationFailedError,
    generate_deck,
)
from ppt_generation.data import dataset_loader, metric_engine


REPO_ROOT = Path(__file__).resolve().parent.parent
BACKEND_ROOT = str(REPO_ROOT / "src" / "backend")


def _workbook(path: Path, sheet: str, headers: list[str], rows: list[list[object]]) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = sheet
    worksheet.append(headers)
    for row in rows:
        worksheet.append(row)
    workbook.save(path)
    workbook.close()
    return path


def _domain_workbook(tmp_path: Path, domain: str) -> Path:
    periods = [f"2026-{month:02d}" for month in range(1, 7)]
    if domain == "food":
        return _workbook(
            tmp_path / "food.xlsx",
            "餐飲營運",
            ["月份", "訂單數", "營收", "評分"],
            [[period, 100 + i * 8, 50000 + i * 3200, 4.1 + i * 0.05] for i, period in enumerate(periods)],
        )
    if domain == "travel":
        return _workbook(
            tmp_path / "travel.xlsx",
            "旅遊表現",
            ["月份", "旅客人次", "旅遊消費", "滿意度"],
            [[period, 800 + i * 35, 120000 + i * 4500, 82 + i] for i, period in enumerate(periods)],
        )
    return _workbook(
        tmp_path / "stocks.xlsx",
        "股票行情",
        ["日期", "收盤價", "成交量"],
        [[period, 100 + i * 1.5, 10000 + i * 500] for i, period in enumerate(periods)],
    )


def _ingest(source: Path, destination: Path) -> tuple[Path, dict]:
    if BACKEND_ROOT not in sys.path:
        sys.path.insert(0, BACKEND_ROOT)
    from app.ingestion.generation_bridge import ingest_excel, save_payload

    payload = ingest_excel(source)
    return save_payload(payload, destination), payload


def _base_metric(store, text: str):
    return next(
        metric
        for metric in store.metrics.values()
        if metric.semantic == "value" and text in metric.name
    )


def test_measure_semantics_block_invalid_cross_domain_share(tmp_path: Path) -> None:
    food = _domain_workbook(tmp_path, "food")
    _, payload = _ingest(food, tmp_path / "food.json")
    loaded = dataset_loader.load_ingestion_result(payload)
    store, report = metric_engine.build_metric_store(loaded)

    rating = _base_metric(store, "評分")
    orders = _base_metric(store, "訂單數")

    assert rating.value_semantic == "score"
    assert rating.aggregation_semantic == "average"
    assert store.get(
        rating.metric_key.replace(".value", ".share"),
        require_computable=False,
    ).computable is False
    assert orders.value_semantic == "count"
    assert orders.aggregation_semantic == "sum"
    assert report.dataset_profiles


def test_stock_price_is_not_treated_as_additive_share(tmp_path: Path) -> None:
    stocks = _domain_workbook(tmp_path, "stocks")
    _, payload = _ingest(stocks, tmp_path / "stocks.json")
    store, _ = metric_engine.build_metric_store(
        dataset_loader.load_ingestion_result(payload)
    )

    close = _base_metric(store, "收盤價")
    assert close.value_semantic == "price"
    assert close.aggregation_semantic == "latest"
    assert store.get(
        close.metric_key.replace(".value", ".share"),
        require_computable=False,
    ).computable is False
    assert store.get(
        close.metric_key.replace(".value", ".period_growth")
    ).computable is True


@pytest.mark.parametrize("domain", ["food", "travel", "stocks"])
def test_cross_domain_full_pipeline_is_native_and_t1_verified(
    tmp_path: Path,
    domain: str,
) -> None:
    source = _domain_workbook(tmp_path, domain)
    ingestion, _ = _ingest(source, tmp_path / f"{domain}.json")
    output_dir = tmp_path / f"out-{domain}"
    request = GenerationRequest(
        job_id=f"cross-domain-{domain}",
        prompt="依資料製作管理層簡報，提出可追溯的觀察與下一步建議。",
        ingestion_path=str(ingestion),
        output_dir=str(output_dir),
        sections=["核心概況", "趨勢與變化", "建議與下一步"],
        deck_title=f"{domain} 資料洞察",
        options={
            "policy": "required",
            "deadline_seconds": 60,
            "render_reserve_seconds": 10,
            "use_fake_llm": True,
            "skip_semantic_review": False,
        },
    )
    result = generate_deck(request.model_dump(mode="json"))

    actual_filenames = {item.filename for item in result.artifacts}
    assert any(name.endswith(".pptx") for name in actual_filenames)
    assert any(name.endswith(".xlsx") for name in actual_filenames)
    assert "verification.json" in actual_filenames
    assert "generation_manifest.json" in actual_filenames
    assert result.verification_passed is True
    assert result.external_checked == result.series_checked > 0

    pptx_artifact = next(item for item in result.artifacts if item.filename.endswith(".pptx"))
    presentation = Presentation(Path(pptx_artifact.path))
    visuals = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_chart or getattr(shape, "has_table", False)
    ]
    assert visuals
    assert all(
        shape.has_chart or getattr(shape, "has_table", False)
        for shape in visuals
    )

    deckspec = json.loads((output_dir / "deckspec.json").read_text(encoding="utf-8"))
    assert deckspec["contract_version"] == "1.0"
    assert deckspec["pages"]
    if domain == "stocks":
        narrative = json.dumps(deckspec["pages"], ensure_ascii=False)
        assert "建議買進" not in narrative
        assert "建議賣出" not in narrative
        assert "保證報酬" not in narrative


def test_required_expired_llm_budget_uses_only_legal_fallbacks(tmp_path: Path) -> None:
    from datetime import datetime, timedelta, timezone

    source = _domain_workbook(tmp_path, "food")
    ingestion, _ = _ingest(source, tmp_path / "required.json")
    output_dir = tmp_path / "required-output"
    request = GenerationRequest(
        job_id="required-expired-budget",
        prompt="依資料產出可追溯簡報",
        ingestion_path=str(ingestion),
        output_dir=str(output_dir),
        source_objects=[
            {
                "bucket": "test-bucket",
                "key": "uploads/required-expired-budget/food.xlsx",
                "filename": source.name,
                "size_bytes": source.stat().st_size,
            }
        ],
        options={
            "policy": "required",
            "deadline_seconds": 30,
            "render_reserve_seconds": 5,
            "use_fake_llm": True,
            "skip_semantic_review": False,
        },
        deadline_at_utc=datetime.now(timezone.utc) - timedelta(seconds=1),
    )
    result = generate_deck(request.model_dump(mode="json"))

    manifest = json.loads(
        (output_dir / "generation_manifest.json").read_text(encoding="utf-8")
    )
    review = json.loads(
        (output_dir / "stages" / "05_review.json").read_text(encoding="utf-8")
    )
    assert manifest["request"]["llm_budget_exhausted"] is True
    assert manifest["input"]["source_objects"] == [
        {
            "bucket": "test-bucket",
            "key": "uploads/required-expired-budget/food.xlsx",
            "filename": source.name,
            "size_bytes": source.stat().st_size,
            "etag": None,
            "sha256": None,
        }
    ]
    assert manifest["contracts"]["deckspec_sha256"] == hashlib.sha256(
        (output_dir / "deckspec.json").read_bytes()
    ).hexdigest()
    assert manifest["llm"]["calls_by_stage"] == {
        "sections": 0,
        "charts": 0,
        "narratives": 0,
        "review": 0,
    }
    assert result.external_checked == result.series_checked > 0
    assert all(item["hard_rules_passed"] is True for item in review["reviews"])
    assert all(
        item["selected_candidate_hard_issues"] == []
        for item in review["reviews"]
    )


def test_measure_profile_prioritizes_column_over_dataset_name() -> None:
    from ppt_generation.data.data_profile import infer_measure_profile

    satisfaction = infer_measure_profile(
        key="satisfaction",
        label="滿意度",
        unit="分",
        dataset_name="營收與滿意度",
    )
    orders = infer_measure_profile(
        key="orders",
        label="訂單數",
        unit="筆",
        dataset_name="轉換率與訂單",
    )

    ambiguous = infer_measure_profile(
        key="value",
        label="數值",
        unit=None,
        dataset_name="股票收盤價",
    )

    assert satisfaction.value_semantic == "score"
    assert orders.value_semantic == "count"
    assert ambiguous.value_semantic == "unknown"
    assert "share" not in satisfaction.allowed_derivations
    assert "share" in orders.allowed_derivations


@pytest.mark.parametrize(
    "advice",
    ["應加碼", "建議持有", "逢低布局", "設定停損", "目標價"],
)
def test_price_rule_blocks_trading_advice_variants(advice: str) -> None:
    from ppt_generation.agents.narrative_writer import PageNarrative
    from ppt_generation.agents.reviewer import run_rule_layer
    from ppt_generation.charts.chart_planner import ChartPlan, resolve_chart_plan
    from ppt_generation.core import placeholders
    from ppt_generation.data.metric_store import MetricSeries, MetricStore

    store = MetricStore(source_files=["stocks.xlsx"])
    store.add(
        MetricSeries(
            metric_key="stocks.price",
            name="收盤價",
            categories=["2026-01", "2026-02"],
            series={"收盤價": [100.0, 101.0]},
            unit="元",
            value_semantic="price",
            aggregation_semantic="latest",
            allowed_derivations=["period_growth", "rank", "forecast"],
            axis_kind="temporal",
        )
    )
    chart = resolve_chart_plan(
        ChartPlan(
            slide_title="價格走勢",
            chart_type="line",
            chart_title="價格走勢",
            metric_key="stocks.price",
            series_names=["收盤價"],
            page_number=1,
        ),
        store,
    )
    narrative = PageNarrative(
        page_number=1,
        slide_title="價格走勢",
        headline=f"{advice}，最新觀察值為 {{{{stocks.price|收盤價|latest}}}}",
        bullets=[],
    )
    narrative.cited_metric_keys = placeholders.cited_metric_keys(
        narrative.all_text
    )

    assert any(
        "不得產生買賣或保證報酬建議" in issue
        for issue in run_rule_layer(narrative, chart, store)
    )


def test_derived_price_metrics_preserve_investment_safety_semantic() -> None:
    from ppt_generation.agents.narrative_writer import PageNarrative
    from ppt_generation.agents.reviewer import run_rule_layer
    from ppt_generation.charts.chart_planner import ChartPlan, resolve_chart_plan
    from ppt_generation.core import placeholders
    from ppt_generation.data.metric_engine import derive_forecast, derive_period_growth
    from ppt_generation.data.metric_store import MetricSeries, MetricStore

    base = MetricSeries(
        metric_key="stocks.price.value",
        name="收盤價",
        categories=["2026-01", "2026-02", "2026-03", "2026-04"],
        series={"收盤價": [100.0, 101.0, 102.0, 103.0]},
        unit="元",
        value_semantic="price",
        aggregation_semantic="latest",
        allowed_derivations=["period_growth", "rank", "forecast"],
        axis_kind="temporal",
    )
    derived_metrics = [derive_period_growth(base), derive_forecast(base)]

    for metric in derived_metrics:
        assert metric.value_semantic == "price"
        store = MetricStore(source_files=["stocks.xlsx"])
        store.add(metric)
        chart = resolve_chart_plan(
            ChartPlan(
                slide_title="價格分析",
                chart_type="line",
                chart_title="價格分析",
                metric_key=metric.metric_key,
                series_names=["收盤價"],
                page_number=1,
            ),
            store,
        )
        narrative = PageNarrative(
            page_number=1,
            slide_title="價格分析",
            headline=(
                "應加碼，最新觀察值為 "
                f"{{{{{metric.metric_key}|收盤價|latest}}}}"
            ),
            bullets=[],
        )
        narrative.cited_metric_keys = placeholders.cited_metric_keys(
            narrative.all_text
        )
        assert any(
            "不得產生買賣或保證報酬建議" in issue
            for issue in run_rule_layer(narrative, chart, store)
        )


def test_pipeline_invocation_contract_rejects_invalid_deadline_reserve() -> None:
    from pydantic import ValidationError

    from ppt_generation.contracts.stages import PipelineRequestContract

    with pytest.raises(ValidationError):
        PipelineRequestContract.model_validate(
            {
                "ingestion_payload": {"contract_version": "1.0"},
                "ingestion_source": "ingestion.json",
                "user_prompt": "make deck",
                "output_dir": "tmp",
                "generation_policy": "required",
                "generation_deadline_seconds": 30,
                "generation_render_reserve_seconds": 30,
            }
        )


def test_generation_fails_closed_for_pending_ingestion_review(
    tmp_path: Path,
) -> None:
    ingestion_path = tmp_path / "pending_ingestion.json"
    ingestion_path.write_text(
        json.dumps(
            {
                "contract_version": "1.0",
                "filename": "ambiguous.xlsx",
                "pipeline_status": "completed_with_warnings",
                "source_files": ["ambiguous.xlsx"],
                "datasets": [
                    {
                        "dataset_id": "ambiguous",
                        "name": "Ambiguous",
                        "filename": "ambiguous.xlsx",
                        "source_kind": "excel",
                        "table_kind": "structured_table",
                        "metadata": {},
                        "columns": [
                            {
                                "key": "value",
                                "label": "Value",
                                "index": 0,
                                "data_type": "number",
                            }
                        ],
                        "records": [],
                        "confidence": 0.74,
                        "requires_human_review": True,
                        "review_status": "pending",
                        "warnings": ["layout confidence below gate"],
                    }
                ],
                "warnings": [],
                "errors": [],
            }
        ),
        encoding="utf-8",
    )
    request = GenerationRequest(
        job_id="pending-review",
        prompt="Generate a trend deck",
        ingestion_path=str(ingestion_path),
        output_dir=str(tmp_path / "output"),
        options={
            "deadline_seconds": 10,
            "render_reserve_seconds": 1,
            "use_fake_llm": True,
        },
    )

    with pytest.raises(
        GenerationFailedError,
        match="尚未通過人工確認",
    ):
        generate_deck(request.model_dump(mode="json"))

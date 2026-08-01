"""
原生表格與熱力圖
=================
對應 FR-2.4（表格一律為 PPT 原生 table 物件，禁止文字方塊拼貼）與
FR-2.3 的熱力圖條目（無原生熱力圖類型，以原生表格 + 儲存格底色模擬）。

與 :mod:`chart_builder` 的分工
-------------------------------
圖表走 ``shapes.add_chart()``，有 chart XML 快取與內嵌 workbook 兩份副本；
表格走 ``shapes.add_table()``，**只有儲存格文字一份**——沒有內嵌 workbook，
右鍵也沒有「編輯資料」。這不是缺陷，是 PowerPoint 表格本來的樣子。

但它有個直接後果：表格數字無法靠「①②天生一致」來保證正確。因此本模組
規定表格文字一律由 :func:`format_value` 這一個函式產生，並讓
:mod:`verification.verify_chart_consistency` 反向解析儲存格文字與稽核
Excel 比對（見該模組的表格比對段）。表格的數字正確性靠這條驗證守著。

熱力圖的色階
-------------
色階只是視覺編碼，不得改變任何數值。分母取該表所有數值的極值，
線性插值到兩個端點色之間；缺值不上色（留白），不塗成最小值的顏色——
把缺值畫成「最冷的格子」等於宣稱它是最小值。
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Sequence

from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Emu, Pt

from pptx.oxml.ns import qn as _qn

from ..core import theme
from .chart_builder import ChartSpec


#: 表頭底色與文字色。改用近黑而非原本的深藍——深藍不在本專案的配色裡，
#: 一份簡報同時出現深藍表頭與紅色圖表就是兩套視覺語彙。
HEADER_FILL = theme.CHART_NEUTRAL
HEADER_FONT_COLOR = theme.INVERSE_COLOR

#: 斑馬紋的淺色列。純白／極淺灰交錯，讓讀者的視線不會跨列跑錯行。
BAND_FILL = RGBColor(0xF7, 0xF7, 0xF7)

#: 非斑馬紋列的底色。**必須是明確的白色，不能留空**——留空的儲存格會顯示
#: 表格樣式的預設底色，而 ``add_table()`` 預設套用的
#: 「Medium Style 2 - Accent 1」是以主題 accent1（本模板為藍 4472C4）上色的。
ROW_FILL = RGBColor(0xFF, 0xFF, 0xFF)

#: PowerPoint 內建「無樣式、無格線」表格樣式的 GUID。
#:
#: 為什麼要改：``shapes.add_table()`` 不接受樣式參數，一律套用
#: 「Medium Style 2 - Accent 1」。那個樣式的 ``wholeTbl`` 帶著 accent1 的
#: 淺色調底色，任何沒有明確填色的儲存格都會渲染成淺藍，而且它的
#: ``firstRow`` / ``bandRow`` 還會再疊一層藍。逐格填色能蓋掉底色，
#: 但漏一格就露一格藍——把樣式本身換成無樣式才是根治。
NO_STYLE_NO_GRID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

#: 表格框線：資料列之間的細線，與表頭下緣的品牌色粗線。
BORDER_COLOR = theme.HAIRLINE_COLOR
BORDER_WIDTH_PT = 0.75
HEADER_RULE_COLOR = theme.ACCENT
HEADER_RULE_WIDTH_PT = 1.5

#: 熱力圖色階兩端，取自 theme 的白 → 台新紅漸層。
HEATMAP_LOW = theme.HEATMAP_LOW
HEATMAP_HIGH = theme.HEATMAP_HIGH

#: 底色亮度低於此門檻時，文字改為白色以維持可讀性（無障礙要求）。
_DARK_BACKGROUND_THRESHOLD = theme.DARK_BACKGROUND_THRESHOLD

#: 表格字級上限。實際字級由 ``theme.fit_table_font_size()`` 依列數再往下收。
TITLE_FONT_SIZE = Pt(14)
HEADER_FONT_SIZE = theme.TABLE_HEADER_FONT_SIZE
BODY_FONT_SIZE = theme.TABLE_BODY_FONT_SIZE

#: 一頁能放得下的資料列上限。超過就不是給人看的表了，應改用圖表或先取 Top N。
MAX_TABLE_ROWS = 20


@dataclass
class TableSpec(ChartSpec):
    """
    原生表格規格。

    刻意繼承 :class:`ChartSpec`：``categories`` 為列標籤、``series`` 為欄，
    形狀與圖表完全相同，稽核 Excel 匯出與 planner 查表都不必分支。
    ``chart_type`` 欄位對表格無意義，不使用。
    """

    #: True 時依數值大小為儲存格上色（熱力圖）。
    heatmap: bool = False
    #: 列標籤欄的表頭文字。
    row_header: str = "項目"
    #: 數值單位，用於格式化與表頭標注。
    unit: str | None = None
    #: 需要以粗體強調的列標籤（例如合計列）。
    emphasize_rows: tuple[str, ...] = field(default_factory=tuple)


class TableBuildError(ValueError):
    """表格規格不可用。"""


# ---------------------------------------------------------------------------
# 數值格式化（唯一入口）
# ---------------------------------------------------------------------------
def format_value(value: float | None, unit: str | None = None) -> str:
    """
    把數值格式化成儲存格文字。

    **表格文字一律經由此函式產生**，因為驗證要反向解析它。任何地方
    自己 f-string 出一個數字，就等於在系統裡開了第二個數字來源。

    缺值輸出 ``"—"``（破折號）而非 ``0`` 或空字串：0 是一個觀測值，
    空字串看起來像漏排版，破折號才明確表示「此格無資料」。
    """
    if value is None:
        return "—"

    if unit in {"%", "％"}:
        return f"{value:.2f}%"

    if unit == "名":
        return f"{int(round(value))}"

    if float(value).is_integer() and abs(value) < 1e15:
        return f"{int(value):,}"

    return f"{value:,.2f}"


def parse_value(text: str, unit: str | None = None) -> float | None:
    """
    :func:`format_value` 的反向操作，供驗證比對使用。

    解析不出來時回傳 None，由呼叫端決定要不要視為失敗——這裡不擅自
    當成 0。
    """
    cleaned = str(text).strip().replace(",", "").replace("％", "").rstrip("%")

    if not cleaned or cleaned == "—":
        return None

    try:
        return float(cleaned)
    except ValueError:
        return None


# ---------------------------------------------------------------------------
# 色階
# ---------------------------------------------------------------------------
#: 色階混色與亮度判斷統一由 theme 提供，圖表與表格共用同一套規則。
_blend = theme.blend
_relative_luminance = theme.relative_luminance


def heatmap_color(
    value: float | None,
    minimum: float,
    maximum: float,
) -> RGBColor | None:
    """
    依數值在 [minimum, maximum] 中的位置取色。

    缺值回傳 None（不上色）。全表同值時一律取低值色——此時色階沒有
    資訊量，用深色會誤導讀者以為每格都是高點。
    """
    if value is None:
        return None

    if maximum <= minimum:
        return HEATMAP_LOW

    return _blend(HEATMAP_LOW, HEATMAP_HIGH, (value - minimum) / (maximum - minimum))


def _value_range(spec: TableSpec) -> tuple[float, float]:
    values = [
        value
        for series in spec.series.values()
        for value in series
        if value is not None
    ]

    if not values:
        raise TableBuildError(f"表格 {spec.title!r} 沒有任何可用數值")

    return min(values), max(values)


# ---------------------------------------------------------------------------
# 表格生成
# ---------------------------------------------------------------------------
def add_native_table(
    slide: Slide,
    spec: TableSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    插入 PPT 原生表格（``shapes.add_table``）。

    ``spec.heatmap`` 為 True 時依數值為儲存格上色，模擬熱力圖。

    Returns:
        ``pptx.table.Table``。
    """
    if not spec.categories:
        raise TableBuildError(f"表格 {spec.title!r} 沒有任何資料列")

    if not spec.series:
        raise TableBuildError(f"表格 {spec.title!r} 沒有任何資料欄")

    if len(spec.categories) > MAX_TABLE_ROWS:
        raise TableBuildError(
            f"表格 {spec.title!r} 有 {len(spec.categories)} 列，"
            f"超過單頁可閱讀上限 {MAX_TABLE_ROWS} 列。"
            "請先取 Top N，或改用圖表呈現。"
        )

    column_names = list(spec.series)
    row_count = len(spec.categories) + 1
    column_count = len(column_names) + 1

    graphic_frame = slide.shapes.add_table(
        row_count, column_count, left, top, width, height
    )
    table = graphic_frame.table

    minimum, maximum = _value_range(spec) if spec.heatmap else (0.0, 0.0)

    # 字級依列數推算。表格與文字框不同：PowerPoint 不會自動縮小表格文字，
    # 列高只會被內容一路撐出投影片外，所以這裡必須主動算。
    body_size = theme.fit_table_font_size(row_count, column_count, int(height))
    header_size = body_size

    # 先把預設的藍色表格樣式換掉，再逐格填色。兩者缺一不可：只換樣式會
    # 讓表格變成純白無層次，只填色則會在漏掉的格子露出藍底。
    _reset_table_style(table)

    _write_header(table, spec, column_names, size=header_size)

    for row_offset, category in enumerate(spec.categories, start=1):
        emphasize = category in spec.emphasize_rows
        # 熱力圖的數值格自己會上色，斑馬紋只會與色階打架；但列標籤欄
        # 仍需要明確底色，否則露出樣式預設的藍。
        if spec.heatmap:
            band = ROW_FILL
        else:
            band = BAND_FILL if row_offset % 2 == 0 else ROW_FILL

        _write_label_cell(
            table.cell(row_offset, 0),
            category,
            emphasize,
            size=body_size,
            band=band,
        )

        for column_offset, name in enumerate(column_names, start=1):
            values = spec.series[name]
            value = values[row_offset - 1] if row_offset - 1 < len(values) else None

            _write_value_cell(
                table.cell(row_offset, column_offset),
                value,
                spec,
                emphasize=emphasize,
                heat=(
                    heatmap_color(value, minimum, maximum)
                    if spec.heatmap
                    else None
                ),
                size=body_size,
                band=band,
            )

    return table


#: ``a:tcPr`` 中框線元素的 schema 順序。框線必須排在填色之前，
#: 否則 PowerPoint 判定檔案需要修復。
_BORDER_TAGS = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")


def _reset_table_style(table) -> None:
    """
    把表格改成無樣式、無格線，並關掉表頭與斑馬紋的樣式強調。

    這是「圖表還是藍色的」的根因修復：不做這件事，任何沒被逐格填色的
    儲存格都會露出 accent1 的藍。
    """
    table.first_row = False
    table.horz_banding = False

    tbl = table._tbl
    tbl_pr = tbl.find(_qn("a:tblPr"))

    if tbl_pr is None:
        tbl_pr = tbl.makeelement(_qn("a:tblPr"), {})
        tbl.insert(0, tbl_pr)

    existing = tbl_pr.find(_qn("a:tableStyleId"))

    if existing is not None:
        tbl_pr.remove(existing)

    # a:tableStyleId 是 a:tblPr 的最後一個子元素，直接 append 即符合 schema。
    style_id = tbl_pr.makeelement(_qn("a:tableStyleId"), {})
    style_id.text = NO_STYLE_NO_GRID
    tbl_pr.append(style_id)


def _set_cell_border(
    cell,
    edge: str,
    color: RGBColor,
    width_pt: float,
) -> None:
    """
    為儲存格的某一邊加框線。

    ``edge`` 為 ``L`` / ``R`` / ``T`` / ``B``。python-pptx 沒有框線的高階
    API，但這裡只寫 ``a:tcPr`` 的子元素，不影響表格內容或任何數值。

    換成無樣式表格後格線也一起沒了，所以框線必須自己畫——沒有任何分隔線
    的數字表，讀者的視線會在列之間跑錯行。
    """
    tag = _qn(f"a:ln{edge}")
    tc_pr = cell._tc.get_or_add_tcPr()

    existing = tc_pr.find(tag)

    if existing is not None:
        tc_pr.remove(existing)

    line = tc_pr.makeelement(
        tag,
        {
            "w": str(int(Pt(width_pt))),
            "cap": "flat",
            "cmpd": "sng",
            "algn": "ctr",
        },
    )
    solid_fill = line.makeelement(_qn("a:solidFill"), {})
    srgb = line.makeelement(_qn("a:srgbClr"), {"val": f"{color}"})
    solid_fill.append(srgb)
    line.append(solid_fill)

    # 插在後續框線標籤與填色之前，維持 a:tcPr 的 schema 順序。
    index = _BORDER_TAGS.index(f"a:ln{edge}")
    later = {_qn(name) for name in _BORDER_TAGS[index + 1 :]}
    border_tags = {_qn(name) for name in _BORDER_TAGS}
    anchor = None

    for child in tc_pr:
        if child.tag in later or child.tag not in border_tags:
            anchor = child
            break

    if anchor is None:
        tc_pr.append(line)
    else:
        anchor.addprevious(line)


def _fill_cell(cell, color: RGBColor | None) -> None:
    if color is None:
        return

    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _write_header(
    table,
    spec: TableSpec,
    column_names: Sequence[str],
    *,
    size,
) -> None:
    headers = [spec.row_header, *column_names]

    for index, text in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = str(text)
        # 框線先寫、填色後寫，才符合 a:tcPr 的 schema 順序。
        _set_cell_border(cell, "B", HEADER_RULE_COLOR, HEADER_RULE_WIDTH_PT)
        _fill_cell(cell, HEADER_FILL)

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                theme.apply_font(
                    run,
                    size=size,
                    bold=True,
                    color=HEADER_FONT_COLOR,
                )


def _write_label_cell(
    cell,
    text: str,
    emphasize: bool,
    *,
    size,
    band: RGBColor | None = None,
) -> None:
    cell.text = str(text)
    _set_cell_border(cell, "B", BORDER_COLOR, BORDER_WIDTH_PT)
    _fill_cell(cell, band)

    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            theme.apply_font(
                run,
                size=size,
                bold=emphasize,
                color=theme.TITLE_COLOR,
            )


def _write_value_cell(
    cell,
    value: float | None,
    spec: TableSpec,
    *,
    emphasize: bool,
    heat: RGBColor | None,
    size,
    band: RGBColor | None = None,
) -> None:
    cell.text = format_value(value, spec.unit)
    _set_cell_border(cell, "B", BORDER_COLOR, BORDER_WIDTH_PT)
    # 熱力圖缺值格沒有色階色，仍要填白——留空會露出樣式預設的藍。
    _fill_cell(cell, heat if heat is not None else band)

    font_color = (
        theme.readable_text_color(heat)
        if heat is not None
        else theme.TITLE_COLOR
    )

    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            theme.apply_font(
                run,
                size=size,
                bold=emphasize,
                color=font_color,
            )


def add_heatmap_table(
    slide: Slide,
    spec: TableSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    熱力圖：原生表格 + 資料驅動的儲存格底色。

    這**不是**真正的圖表物件，使用者右鍵沒有「編輯資料」選項——
    PowerPoint 沒有熱力圖類型，這是規格書 FR-2.3 明訂的替代方案。
    """
    heat_spec = TableSpec(
        title=spec.title,
        categories=spec.categories,
        series=spec.series,
        heatmap=True,
        row_header=spec.row_header,
        unit=spec.unit,
        emphasize_rows=spec.emphasize_rows,
    )

    return add_native_table(slide, heat_spec, left, top, width, height)


#: 表格類 skill registry。與 CHART_SKILLS 分開註冊，因為它們的驗收標準不同：
#: 圖表要能右鍵編輯資料，表格不能（也不該假裝可以）。
TABLE_SKILLS = {
    "table": add_native_table,
    "heatmap": add_heatmap_table,
}


TABLE_SKILL_TOOL_SCHEMAS = [
    {
        "name": "table",
        "description": (
            "PPT 原生表格。適合需要精確讀出每一格數字的場合"
            "（如風險指標對照、Top N 明細）。列數上限 "
            f"{MAX_TABLE_ROWS}，超過請先取 Top N。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "description": "必填；只列出與本頁主題直接相關的系列。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "heatmap",
        "description": (
            "熱力圖：原生表格 + 資料驅動的儲存格底色。"
            "適合多實體 × 多期間的強弱分佈一眼比較。"
            "注意這不是圖表物件，右鍵沒有「編輯資料」。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "description": "必填；只列出與本頁主題直接相關的系列。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
]
"""給 LLM 的表格 skill schema。與圖表一樣沒有數值欄位。"""

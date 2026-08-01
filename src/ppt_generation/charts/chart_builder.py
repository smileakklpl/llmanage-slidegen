"""
圖表生成模組 (POC)
=====================
核心設計原則：
1. 圖表數字的「唯一真相來源」是本模組的輸入資料（解析後的 Excel 結構化資料）。
2. 一律透過 python-pptx 的 add_chart() API 生成圖表，
   讓「畫面顯示值 (chart cache)」與「右鍵編輯資料看到的內嵌工作表」
   自動保持一致，不手動操作底層 XML。
3. 不使用外部連結 (external link) 方式綁定 Excel，
   因為 PPT 寄出後外部路徑會失效。改用 PowerPoint 原生的
   「內嵌工作簿 (embedded workbook)」機制，這正是使用者
   右鍵「編輯資料」時開啟的資料來源。
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Sequence

from pptx.chart.axis import CategoryAxis, ValueAxis
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    XL_TICK_MARK,
)
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from lxml import etree

from ..core import theme


#: 這些圖型沒有座標軸，套用軸樣式會拋 ValueError。
_AXISLESS_CHART_TYPES = frozenset(
    {
        XL_CHART_TYPE.PIE,
        XL_CHART_TYPE.PIE_EXPLODED,
        XL_CHART_TYPE.DOUGHNUT,
        XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    }
)

#: 折線類圖型的顏色要設在線上（``a:ln``），不是填滿。
_LINE_CHART_TYPES = frozenset(
    {
        XL_CHART_TYPE.LINE,
        XL_CHART_TYPE.LINE_MARKERS,
        XL_CHART_TYPE.LINE_STACKED,
        XL_CHART_TYPE.LINE_MARKERS_STACKED,
        XL_CHART_TYPE.XY_SCATTER,
        XL_CHART_TYPE.XY_SCATTER_LINES,
    }
)


@dataclass
class ChartSpec:
    """一張圖表所需的資料規格。"""

    title: str
    categories: Sequence[str]
    series: dict[str, Sequence[float]]  # {系列名稱: 數值}
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED


# ---------------------------------------------------------------------------
# 視覺樣式
# ---------------------------------------------------------------------------
# 樣式只碰 spPr（形狀屬性）與 txPr（文字屬性），完全不觸碰 c:numCache /
# c:strCache 等數值節點，因此「chart XML 快取 ↔ 內嵌 workbook ↔ 稽核 xlsx」
# 三份副本的一致性不受影響——那是 add_chart() 一次寫好的，與上色無關。
# ---------------------------------------------------------------------------
def _style_chart_area(chart) -> None:
    """
    圖表底板與繪圖區改為透明、無外框。

    不設的話兩者都沒有 ``c:spPr``，PowerPoint 會套用主題預設——底板可能
    帶上一圈框線或淺色底，而那個顏色來自主題的 accent／bg，本專案控制不到。
    投影片本身是白底，圖表底板透明才是對的。
    """
    chart_space = chart._chartSpace
    chart_element = chart_space.find(qn("c:chart"))

    # CT_ChartSpace 的順序：… c:chart → c:spPr → c:txPr → c:externalData。
    # 所以底板的 spPr 要緊接在 c:chart 之後，不能 append 到最後
    # （那會排到 externalData 後面，PowerPoint 判定檔案損毀）。
    existing = chart_space.find(qn("c:spPr"))

    if existing is not None:
        chart_space.remove(existing)

    chart_element.addnext(_transparent_sp_pr(chart_space))

    plot_area = chart_element.find(qn("c:plotArea"))

    if plot_area is None:
        return

    # CT_PlotArea 的 c:spPr 是最後一個子元素。雙軸圖會在套樣式之後再追加
    # 次軸，所以這裡先移除舊的再重新 append，確保它永遠在最後。
    existing_plot_sp_pr = plot_area.find(qn("c:spPr"))

    if existing_plot_sp_pr is not None:
        plot_area.remove(existing_plot_sp_pr)

    plot_area.append(_transparent_sp_pr(plot_area))


def _transparent_sp_pr(parent):
    """建立一個「無填色、無框線」的 ``c:spPr``。"""
    sp_pr = parent.makeelement(qn("c:spPr"), {})
    sp_pr.append(sp_pr.makeelement(qn("a:noFill"), {}))
    line = sp_pr.makeelement(qn("a:ln"), {})
    line.append(sp_pr.makeelement(qn("a:noFill"), {}))
    sp_pr.append(line)

    return sp_pr


def _style_chart_title(chart, text: str) -> None:
    """圖表標題：左上、12pt、深灰。不用模板預設的 18pt 置中大標。"""
    chart.has_title = True
    text_frame = chart.chart_title.text_frame
    text_frame.text = text

    for paragraph in text_frame.paragraphs:
        theme.apply_chart_font(
            paragraph.font,
            size=theme.CHART_TITLE_FONT_SIZE,
            bold=True,
            color=theme.TITLE_COLOR,
        )


def _style_axes(chart) -> None:
    """
    座標軸：淺格線、無刻度線、9pt 標籤。

    刻意直接走 ``catAx_lst`` / ``valAx_lst`` 而不用 ``chart.category_axis``
    與 ``chart.value_axis``：雙軸圖有兩個 ``c:valAx``，而 python-pptx 的
    ``value_axis`` 在有兩軸時回傳的是**次軸**，主軸會被漏掉。
    """
    chart_space = chart._chartSpace

    for element in chart_space.catAx_lst:
        axis = CategoryAxis(element)
        axis.has_major_gridlines = False
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.format.line.color.rgb = theme.HAIRLINE_COLOR
        theme.apply_chart_font(
            axis.tick_labels.font,
            size=theme.CHART_LABEL_FONT_SIZE,
            color=theme.BODY_COLOR,
        )

    for element in chart_space.valAx_lst:
        axis = ValueAxis(element)
        # 橫向格線留著（讀者靠它比高度），但要淺到不與長條爭視線；
        # 軸線本身反而可以省掉，格線已經界定了範圍。
        axis.has_major_gridlines = True
        axis.major_gridlines.format.line.color.rgb = theme.HAIRLINE_COLOR
        axis.major_gridlines.format.line.width = Pt(0.75)
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.format.line.fill.background()
        theme.apply_chart_font(
            axis.tick_labels.font,
            size=theme.CHART_LABEL_FONT_SIZE,
            color=theme.MUTED_COLOR,
        )


def _style_legend(chart, series_count: int) -> None:
    """
    單系列不放圖例，多系列放底部。

    一個系列的圖例只是把系列名稱重複一次，而圖表標題已經寫了同一件事，
    卻要吃掉右側或底部一整條空間。
    """
    if series_count <= 1:
        chart.has_legend = False
        return

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    theme.apply_chart_font(
        chart.legend.font,
        size=theme.CHART_LEGEND_FONT_SIZE,
        color=theme.BODY_COLOR,
    )


def _set_series_color(series, color, *, as_line: bool) -> None:
    """
    為一個系列上色。

    折線系列的顏色在 ``a:ln`` 上，長條／圓餅在填滿上。設錯位置的結果是
    折線仍是模板預設色，而 PowerPoint 又不會報錯，只是圖看起來沒改。
    """
    if as_line:
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.25)
        # 雙軸圖的折線系列是先以長條建立、之後才改造的，spPr 上還留著
        # 當長條時設的 solidFill。留著它會讓 PowerPoint 用填滿色去畫標記，
        # 與線色不一致。
        series.format.fill.background()

        marker = getattr(series, "marker", None)

        if marker is not None:
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = color
            marker.format.line.color.rgb = theme.INVERSE_COLOR

        return

    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    # 長條之間不描邊。描邊在資料點多時會讓整片圖看起來糊掉。
    series.format.line.fill.background()


def _color_points_by_rank(series, values: Sequence[float | None]) -> None:
    """
    依數值排名為單一系列的每個資料點上色（深＝大）。

    單系列圖表（排名圖、市占率圓餅圖）的「系列」只有一個，顏色若只設在
    系列上，整張圖就是一種紅——此時顏色沒有承載任何資訊。改為逐點上色，
    顏色深淺即數值大小。
    """
    # 系列層級也給一個保底色。逐點上色理論上覆蓋每個點，但只要有一個點
    # 漏掉（例如未來改動讓 dPt 數量與資料點對不上），漏掉的那個就會落回
    # PowerPoint 自動配色——在本模板下那是主題 accent1 的藍。保底色讓
    # 最壞情況是「深紅」而不是「一根藍柱」。
    _set_series_color(series, theme.CHART_RAMP_HIGH, as_line=False)

    colors = theme.rank_ramp_colors(list(values))

    for index, color in enumerate(colors):
        if index >= len(series.points):
            break

        point = series.points[index]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color
        point.format.line.color.rgb = theme.INVERSE_COLOR
        point.format.line.width = Pt(0.75)


def apply_chart_style(
    chart,
    spec: ChartSpec,
    *,
    line_series_names: Sequence[str] = (),
) -> None:
    """
    套用單色階（白 → 台新紅 + 近黑）配色與字體。

    這是所有圖表的唯一樣式入口。不套用的結果是圖表繼承模板主題的
    accent1..accent6 預設調色盤——一張圖出現藍、橘、灰、黃四種不相干的
    顏色，讀者無法從顏色讀出任何資訊，也與台新紅的品牌語彙互相打架。

    配色規則：

    - **單系列**：逐資料點依排名取白→紅漸層，深紅即最大值
    - **多系列**：系列間用紅／黑交錯（:func:`theme.series_colors`）

    Args:
        chart: ``add_chart()`` 回傳的圖表物件。
        spec: 產生該圖表的規格，用來取系列名稱與數值。
        line_series_names: 已被改造成折線的系列名稱（雙軸圖用）。
            這些系列的顏色會設在線上而非填滿。
    """
    # 圖表層級的預設字體。python-pptx 建立的 chartSpace 帶著 sz="1800"，
    # 任何沒被逐一設定的元素（軸標題、未來新增的資料標籤）都會以 18pt
    # 渲染，在 60% 寬的圖表區裡直接糊成一團。先把預設值收到 9pt，
    # 後續各元素再各自覆寫。
    theme.apply_chart_font(
        chart.font,
        size=theme.CHART_LABEL_FONT_SIZE,
        color=theme.BODY_COLOR,
    )

    _style_chart_area(chart)

    _style_chart_title(chart, spec.title)

    if spec.chart_type not in _AXISLESS_CHART_TYPES:
        _style_axes(chart)

    series_names = list(spec.series)
    _style_legend(chart, len(series_names))

    is_line_chart = spec.chart_type in _LINE_CHART_TYPES
    palette = theme.series_colors(len(series_names))

    # chart.series 的順序與 spec.series 一致（雙軸圖改造後，被搬到
    # lineChart 的系列會排在後面，因此改用名稱查表而非位置對應）。
    series_by_name = {series.name: series for series in chart.series}

    for index, name in enumerate(series_names):
        series = series_by_name.get(name)

        if series is None:
            continue

        as_line = is_line_chart or name in line_series_names

        if len(series_names) == 1 and not as_line:
            # 單系列：顏色留給資料點承載數值大小。
            _color_points_by_rank(series, spec.series[name])
        else:
            _set_series_color(series, palette[index], as_line=as_line)


def _style_pie_labels(chart) -> None:
    """圓餅圖標上百分比。看圓餅估比例本來就不準，數字要直接寫上去。"""
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.show_percentage = True
    labels.show_value = False
    labels.show_series_name = False
    labels.number_format = "0.0%"
    labels.number_format_is_linked = False
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    theme.apply_chart_font(
        labels.font,
        size=theme.CHART_DATA_LABEL_FONT_SIZE,
        color=theme.BODY_COLOR,
    )


def add_category_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    將 ChartSpec 的資料寫入圖表，並插入到指定投影片。

    重點：CategoryChartData 是唯一資料入口。
    python-pptx 會同時：
      1. 寫入 chart1.xml 的 <c:numCache>（畫面顯示用）
      2. 產生對應的 embedded .xlsx（右鍵「編輯資料」開啟用）
    兩者保證數值一致，因為都是從這裡的同一份 spec 產生。

    資料寫入後套用 :func:`apply_chart_style`，只改 spPr／txPr，
    不觸碰任何數值節點。
    """
    chart_data = CategoryChartData()
    chart_data.categories = list(spec.categories)
    for series_name, values in spec.series.items():
        chart_data.add_series(series_name, values)

    graphic_frame = slide.shapes.add_chart(
        spec.chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    apply_chart_style(chart, spec)
    return chart


@dataclass
class ScatterSpec:
    """散點圖（規模 vs 成長）專用規格。"""

    title: str
    series_name: str
    points: Sequence[tuple[float, float]]  # [(x, y), ...]
    labels: Sequence[str] | None = None  # 對應每個點的標籤（如銀行名稱）


def add_scatter_chart(
    slide: Slide,
    spec: ScatterSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    散點圖走 XyChartData，一樣是唯一資料入口原則。
    labels 部分 python-pptx 原生不支援資料點標籤文字，
    若模板要求顯示銀行名稱標籤，需要在 chart 生成後
    透過底層 XML（c:dLbls）額外處理，這裡先留 TODO。
    """
    chart_data = XyChartData()
    series = chart_data.add_series(spec.series_name)
    for x, y in spec.points:
        series.add_data_point(x, y)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart

    theme.apply_chart_font(
        chart.font,
        size=theme.CHART_LABEL_FONT_SIZE,
        color=theme.BODY_COLOR,
    )
    _style_chart_title(chart, spec.title)
    _style_axes(chart)
    # 散點圖只有一個系列，圖例只是把系列名重複一次。
    chart.has_legend = False

    # 散點圖的顏色設在標記上，不是線上——XY_SCATTER 沒有連線。
    for chart_series in chart.series:
        chart_series.format.line.fill.background()
        chart_series.marker.style = XL_MARKER_STYLE.CIRCLE
        chart_series.marker.size = 8
        chart_series.marker.format.fill.solid()
        chart_series.marker.format.fill.fore_color.rgb = theme.ACCENT
        chart_series.marker.format.line.color.rgb = theme.INVERSE_COLOR

    # TODO: 若需散點標籤（如銀行名稱），需手動操作 chart.plots[0] 的
    # dLbls XML 節點，python-pptx 目前無高階 API。
    return chart


def add_pie_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(6096000),
    height: int = Emu(4114800),
):
    """
    市占率圖（圓餅圖）。走同一個 CategoryChartData 入口，
    僅 chart_type 固定為 PIE。要求 spec.series 只能有一組系列。
    """
    if len(spec.series) != 1:
        raise ValueError("圓餅圖只能有一組系列，請確認 ChartSpec.series 長度為 1")
    pie_spec = ChartSpec(
        title=spec.title,
        categories=spec.categories,
        series=spec.series,
        chart_type=XL_CHART_TYPE.PIE,
    )
    chart = add_category_chart(slide, pie_spec, left, top, width, height)

    # 圓餅圖的圖例是必要的：切片上只標百分比，類別名稱得靠圖例對照。
    # 這是 _style_legend 的單系列規則唯一該被覆寫的場合。
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    theme.apply_chart_font(
        chart.legend.font,
        size=theme.CHART_LEGEND_FONT_SIZE,
        color=theme.BODY_COLOR,
    )

    _style_pie_labels(chart)

    return chart


# ---------------------------------------------------------------------------
# 雙軸圖（長條 + 折線，次要數值軸）
# ---------------------------------------------------------------------------
# 這是 python-pptx 唯一沒有高階 API 的必做圖型（附件三 P.5/P.6 官方主打頁：
# 流通卡數長條 + 簽帳金額折線，兩者量級差一個數量級，非用次軸不可）。
#
# 做法上有兩條路，只有一條是對的：
#
# ✗ 自己組 chart XML 並手填 <c:numCache> —— 快取值會有，但內嵌 workbook 不會有，
#   使用者右鍵「編輯資料」看到空表，或看到與畫面不符的舊值。這正是附件三的病灶。
#
# ✓ 先用 add_chart() 一次把**所有**系列（含之後要變折線的）寫進去，
#   內嵌 workbook 由 python-pptx 完整寫好；再把後段 <c:ser> 節點從
#   <c:barChart> 搬到新建的 <c:lineChart> 並掛上次軸。
#
# 搬動的是「這個系列用什麼圖形畫」，完全沒有碰任何數值節點，
# 所以三份副本的一致性由 add_chart() 保證，與圖形改造無關。
@dataclass
class ComboSpec(ChartSpec):
    """
    雙軸圖規格。

    ``series`` 的順序即繪製順序；``line_series_names`` 列出的系列會改用折線
    並掛到次要數值軸，其餘留在主軸長條。刻意繼承 :class:`ChartSpec`——
    稽核 Excel 匯出與三方比對都只認 ChartSpec 的形狀，繼承讓它們不必分支。
    """

    line_series_names: tuple[str, ...] = field(default_factory=tuple)

    @property
    def column_series_names(self) -> tuple[str, ...]:
        return tuple(
            name for name in self.series if name not in self.line_series_names
        )


def _next_axis_ids(plot_area, count: int = 2) -> list[int]:
    """產生未被使用的 axId。與現有 id 相撞會讓 PowerPoint 判定檔案損毀。"""
    existing = {
        int(element.get("val"))
        for element in plot_area.iter(qn("c:axId"))
        if element.get("val", "").isdigit()
    }
    candidate = (max(existing) if existing else 100_000_000) + 1

    ids: list[int] = []

    while len(ids) < count:
        if candidate not in existing:
            ids.append(candidate)

        candidate += 1

    return ids


def _sub(parent, tag: str, **attrs):
    element = etree.SubElement(parent, qn(tag))

    for name, value in attrs.items():
        element.set(name, str(value))

    return element


def _append_secondary_axes(plot_area, cat_ax_id: int, val_ax_id: int) -> None:
    """
    加上次軸的 catAx / valAx 配對。

    次軸的類別軸一定要 ``delete=1``：兩組類別軸都畫出來的話，圖底下會出現
    兩排一模一樣的月份標籤。數值軸則放右側（``axPos=r``）。
    """
    # 座標軸在 CT_PlotArea 中排在 c:dTable / c:spPr 之前。套過樣式的圖表
    # 已經有 c:spPr，直接 append 會讓軸排到它後面而違反 schema。
    trailing = None

    for tag in ("c:dTable", "c:spPr"):
        found = plot_area.find(qn(tag))

        if found is not None:
            trailing = found
            break

    def _new_axis(tag: str):
        element = plot_area.makeelement(qn(tag), {})

        if trailing is None:
            plot_area.append(element)
        else:
            trailing.addprevious(element)

        return element

    cat_ax = _new_axis("c:catAx")
    _sub(cat_ax, "c:axId", val=cat_ax_id)
    _sub(_sub(cat_ax, "c:scaling"), "c:orientation", val="minMax")
    _sub(cat_ax, "c:delete", val=1)
    _sub(cat_ax, "c:axPos", val="b")
    _sub(cat_ax, "c:crossAx", val=val_ax_id)

    val_ax = _new_axis("c:valAx")
    _sub(val_ax, "c:axId", val=val_ax_id)
    _sub(_sub(val_ax, "c:scaling"), "c:orientation", val="minMax")
    _sub(val_ax, "c:delete", val=0)
    _sub(val_ax, "c:axPos", val="r")
    _sub(val_ax, "c:numFmt", formatCode="General", sourceLinked="0")
    _sub(val_ax, "c:majorTickMark", val="out")
    _sub(val_ax, "c:minorTickMark", val="none")
    _sub(val_ax, "c:tickLblPos", val="nextTo")
    _sub(val_ax, "c:crossAx", val=cat_ax_id)


def _promote_series_to_line(chart, line_indices: Sequence[int]) -> None:
    """
    把指定索引的 ``c:ser`` 從 barChart 搬到新建的 lineChart，並掛上次軸。

    ``line_indices`` 是系列在 barChart 中的位置（0 起算），對應 ComboSpec
    中 ``series`` 的順序。
    """
    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))

    if bar_chart is None:
        raise ValueError(
            "找不到 c:barChart，雙軸圖必須以 COLUMN_CLUSTERED 為底圖建立"
        )

    all_series = bar_chart.findall(qn("c:ser"))
    moving = [all_series[index] for index in line_indices]

    if len(moving) == len(all_series):
        raise ValueError("雙軸圖至少要留一個系列在長條主軸上")

    cat_ax_id, val_ax_id = _next_axis_ids(plot_area)

    line_chart = etree.Element(qn("c:lineChart"))
    _sub(line_chart, "c:grouping", val="standard")
    _sub(line_chart, "c:varyColors", val=0)

    for ser in moving:
        bar_chart.remove(ser)

        # invertIfNegative 是長條圖專屬子元素，留在折線 ser 裡會違反
        # DrawingML schema，PowerPoint 會直接判定檔案需要修復。
        for tag in ("c:invertIfNegative",):
            child = ser.find(qn(tag))

            if child is not None:
                ser.remove(child)

        marker = etree.Element(qn("c:marker"))
        _sub(marker, "c:symbol", val="circle")
        _sub(marker, "c:size", val=7)

        # marker 在 CT_LineSer 中的位置在 spPr 之後、cat 之前。
        anchor = ser.find(qn("c:cat"))

        if anchor is not None:
            anchor.addprevious(marker)
        else:
            ser.append(marker)

        line_chart.append(ser)
        _sub(ser, "c:smooth", val=0)

    _sub(line_chart, "c:marker", val=1)
    _sub(line_chart, "c:axId", val=cat_ax_id)
    _sub(line_chart, "c:axId", val=val_ax_id)

    # CT_PlotArea 的元素順序：layout → 各圖形群組 → 各軸 → dTable → spPr。
    # lineChart 必須緊接在 barChart 之後，軸則接在既有軸之後。
    bar_chart.addnext(line_chart)
    _append_secondary_axes(plot_area, cat_ax_id, val_ax_id)


def add_combo_chart(
    slide: Slide,
    spec: ComboSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    雙軸圖：部分系列畫長條（主軸），部分畫折線（次軸）。

    資料入口仍然是 :func:`add_category_chart` 的 ``add_chart()``——
    先讓 python-pptx 把全部系列與內嵌 workbook 寫好，再改變後段系列的
    圖形與所屬軸。全程不觸碰任何數值節點。
    """
    if not spec.line_series_names:
        raise ValueError(
            "雙軸圖需指定 line_series_names（要畫成折線並掛次軸的系列）"
        )

    series_names = list(spec.series)
    unknown = [
        name for name in spec.line_series_names if name not in spec.series
    ]

    if unknown:
        raise ValueError(
            f"line_series_names {unknown} 不在 series 中，"
            f"可用系列：{series_names}"
        )

    if len(spec.line_series_names) >= len(series_names):
        raise ValueError(
            "雙軸圖至少要留一個系列畫長條，"
            f"目前 {len(series_names)} 個系列全被指定為折線"
        )

    base = ChartSpec(
        title=spec.title,
        categories=spec.categories,
        series=spec.series,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
    )
    chart = add_category_chart(slide, base, left, top, width, height)

    _promote_series_to_line(
        chart,
        [
            index
            for index, name in enumerate(series_names)
            if name in spec.line_series_names
        ],
    )

    # 樣式必須在改造**之後**再套一次。改造前這些系列還是長條，顏色設在
    # 填滿上；搬到 lineChart 後填滿不再決定線色，得改設 a:ln。同時次軸
    # 是改造時才新建的，改造前不存在，_style_axes 也就套不到它。
    apply_chart_style(chart, spec, line_series_names=spec.line_series_names)

    return chart


# ---------------------------------------------------------------------------
# Skill Registry：圖表 Agent 只透過 skill 名稱字串呼叫，不需知道實作細節。
# 新增圖表類型時，在此註冊即可，Agent 端無需改動。
# ---------------------------------------------------------------------------
CHART_SKILLS = {
    "column": add_category_chart,   # 長條圖／排名圖／成長率圖
    "bar": add_category_chart,      # 橫條圖
    "line": add_category_chart,     # 折線圖／時間趨勢圖
    "pie": add_pie_chart,           # 市占率圖
    "scatter": add_scatter_chart,   # 規模 vs 成長散點圖
    "combo": add_combo_chart,       # 雙軸圖：長條（主軸）+ 折線（次軸）
    # "heatmap": 由原生表格 + 儲存格底色模擬，不走 add_chart()，
    #            另外註冊在 table_builder.py（尚未實作），Agent 需個別處理。
}

#: skill 名稱 → python-pptx 圖表類型。
#: chart_planner 依此設定 ChartSpec.chart_type，讓 CHART_SKILLS 中共用
#: add_category_chart 的多個 skill（column/bar/line）產生不同圖形。
CHART_TYPE_BY_SKILL = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    # 雙軸圖以長條為底圖建立，折線系列在 add_chart() 之後才改造。
    "combo": XL_CHART_TYPE.COLUMN_CLUSTERED,
}


# ---------------------------------------------------------------------------
# Function Calling / Tool Use 對應層
# ---------------------------------------------------------------------------
# LLM API（OpenAI / Anthropic / Bedrock 等）本身無法執行 Python 檔案，
# 它只能根據我們提供的 "tool schema" 回傳一段結構化 JSON，
# 表達「我想呼叫哪個 skill、帶哪些參數」。
#
# 真正的執行流程是：
#   1. 我們把 CHART_SKILLS 的 key 轉成 tool schema，隨 prompt 送給 LLM
#   2. LLM 回傳 {"name": "pie", "arguments": {...}}（純文字/JSON，不是程式碼）
#   3. 我們自己的 dispatcher（如下 dispatch_chart_skill）收到這段 JSON，
#      在本地 Python 環境查表、驗證、實際呼叫對應函式
#
# LLM 從未執行過任何一行本地程式碼，執行權限完全留在我方。
# ---------------------------------------------------------------------------

CHART_SKILL_TOOL_SCHEMAS = [
    {
        "name": "column",
        "description": "長條圖，適合排名、成長率比較等類別型數據。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "maxItems": 2,
                    "description": "必填；只列出與本頁主題直接相關的系列名稱。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "bar",
        "description": (
            "橫條圖，適合類別名稱較長的排名比較（如銀行名稱橫向排列）。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "maxItems": 2,
                    "description": "必填；只列出與本頁主題直接相關的系列名稱。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "line",
        "description": (
            "折線圖，適合時間序列趨勢。僅可用於 axis_kind 為 temporal 的指標。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "maxItems": 2,
                    "description": "必填；只列出與本頁主題直接相關的系列名稱。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "pie",
        "description": "圓餅圖，適合市占率等單一系列的佔比數據。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {"type": "string"},
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "maxItems": 1,
                    "description": "必填且只能指定一個與本頁主題相關的系列。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "combo",
        "description": (
            "雙軸圖（長條 + 折線）。適合兩個量級差距大、但需並列比較的系列，"
            "例如流通卡數（張）與簽帳金額（百萬元）同時看趨勢。"
            "series_names 的第一個畫長條掛主軸，第二個畫折線掛右側次軸。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 2,
                    "maxItems": 2,
                    "description": (
                        "恰好兩個系列名稱。第一個為長條（主軸），"
                        "第二個為折線（次軸）。"
                    ),
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
    {
        "name": "scatter",
        "description": "散點圖，適合「規模 vs 成長」等雙變數關係。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {"type": "string"},
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 2,
                    "maxItems": 2,
                    "description": "恰好兩個系列名稱，分別對應 x 軸與 y 軸。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
]
"""
給 LLM API 的 tool schema 清單。注意 schema 中完全沒有 "values" 這種
數字欄位，LLM 只能填 metric_key 引用，這是刻意設計，逼迫 LLM 無法
自己編造數字（呼應 .kiro/steering/tech.md 第 1 條原則）。

實際串接時（OpenAI 範例）：

    response = client.chat.completions.create(
        model=...,
        messages=[...],
        tools=[{"type": "function", "function": schema}
               for schema in CHART_SKILL_TOOL_SCHEMAS],
    )
    tool_call = response.choices[0].message.tool_calls[0]
    skill_name = tool_call.function.name          # 例如 "pie"
    raw_args = json.loads(tool_call.function.arguments)

Anthropic / Bedrock 的 tool use 格式略有不同（key 名稱不同），
但核心流程一致：LLM 回傳「工具名稱 + 參數」，執行仍在本地端。
"""


def dispatch_chart_skill(skill_name: str) -> callable:
    """
    Dispatcher：LLM 回傳的 skill_name（純字串）透過這裡查表，
    取得實際的 Python 函式並在本地執行。

    這是 LLM 輸出與本地程式碼執行之間唯一的橋接點，
    也是防呆的關卡：LLM 若回傳未註冊的名稱，直接拋錯，
    不會有機會執行任何非預期的程式碼路徑
    （因為只從 CHART_SKILLS 這個白名單裡取值，不會 eval/exec LLM 的輸出）。
    """
    skill = CHART_SKILLS.get(skill_name)
    if skill is None:
        raise ValueError(
            f"LLM 回傳了未註冊的 skill: {skill_name!r}，"
            f"可用選項: {list(CHART_SKILLS.keys())}"
        )
    return skill

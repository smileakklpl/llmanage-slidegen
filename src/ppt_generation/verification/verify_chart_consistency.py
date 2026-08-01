"""
三方數值比對驗證（對應規格書 T1）
==================================
驗證同一組數字在三個副本中完全一致：

| 副本 | 位置 | 使用者何時看到 |
|---|---|---|
| ① Chart XML 快取 | `.pptx` 內 `ppt/charts/chart#.xml` | 開啟 PPT 看到的圖 |
| ② 內嵌 workbook | `.pptx` 內 `ppt/embeddings/*.xlsx` | 圖表右鍵「編輯資料」 |
| ③ 外部稽核 `.xlsx` | 獨立檔案，隨 PPT 一起寄送 | 獨立核對數字 |

①②由同一次 ``add_chart()`` 呼叫產生；③由 :mod:`excel_exporter` 從相同的
ChartSpec 產生。三者若有任一不符，即為附件三所揭露的「簡報數字與 Excel
不符」問題，必須視為缺陷。

用法::

    python -m ppt_generation.verify_chart_consistency \\
        outputs/deck.pptx outputs/deck_data.xlsx
"""

from __future__ import annotations

import argparse
import io
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Sequence

from openpyxl import load_workbook
from pptx import Presentation

from ..charts import table_builder

#: 浮點比較容許誤差。三份副本理應完全相同，此值只用於吸收
#: Excel 儲存 float 時的最後一位表示差異，不容許實質數值差異。
TOLERANCE = 1e-9

#: 稽核 Excel 中資料表的表頭關鍵字。
_CATEGORY_HEADERS = ("類別", "標籤")


@dataclass
class SeriesComparison:
    """
    單一系列的三方比對結果。

    原生表格沒有內嵌 workbook（PowerPoint 表格本來就沒有「編輯資料」），
    此時 ``embedded`` 為 None，只比對①儲存格文字與③稽核 Excel。
    這不是放寬標準——表格的畫面值就是它唯一的一份值，沒有第二份可比。
    """

    slide_number: int
    chart_title: str
    series_name: str
    chart_cache: list[float | None]
    embedded: list[float | None] | None
    external: list[float | None] | None = None
    #: 這一筆是原生表格（走 add_table）而非原生圖表（走 add_chart）。
    is_table: bool = False

    @property
    def cache_matches_embedded(self) -> bool:
        if self.embedded is None:
            return True

        return _values_equal(self.chart_cache, self.embedded)

    @property
    def cache_matches_external(self) -> bool | None:
        if self.external is None:
            return None

        return _values_equal(self.chart_cache, self.external)

    @property
    def passed(self) -> bool:
        if not self.cache_matches_embedded:
            return False

        # 表格只有一份副本可比。沒有稽核 Excel 就等於這頁的數字沒有任何
        # 獨立來源可以核對——那是「無法驗證」，不是「通過」。圖表不同，
        # 它至少還有①②兩份可互相印證。
        if self.is_table and self.external is None:
            return False

        return self.cache_matches_external is not False

    @property
    def copies_compared(self) -> int:
        """實際比對到的副本數。表格只有兩份（畫面 + 稽核 Excel）。"""
        return 1 + (self.embedded is not None) + (self.external is not None)

    def describe_failure(self) -> str:
        parts: list[str] = []

        if not self.cache_matches_embedded:
            parts.append(
                f"畫面顯示值 {self.chart_cache} 與內嵌工作表 {self.embedded} 不符"
            )

        if self.cache_matches_external is False:
            parts.append(
                f"畫面顯示值 {self.chart_cache} 與外部稽核 Excel "
                f"{self.external} 不符"
            )

        return "；".join(parts)


@dataclass
class VerificationReport:
    """整份簡報的驗證結果。"""

    comparisons: list[SeriesComparison] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return bool(self.comparisons) and all(
            comparison.passed for comparison in self.comparisons
        )

    @property
    def failures(self) -> list[SeriesComparison]:
        return [c for c in self.comparisons if not c.passed]

    @property
    def external_checked(self) -> int:
        return sum(1 for c in self.comparisons if c.external is not None)


def _values_equal(
    left: Sequence[float | None],
    right: Sequence[float | None],
) -> bool:
    """逐項比較，None 只與 None 相等。"""
    if len(left) != len(right):
        return False

    for a, b in zip(left, right):
        if a is None or b is None:
            if a is not b and (a is not None or b is not None):
                return False

            continue

        if abs(float(a) - float(b)) > TOLERANCE:
            return False

    return True


def _to_float(value: Any) -> float | None:
    if value is None:
        return None

    if isinstance(value, bool):
        return None

    if isinstance(value, (int, float)):
        return float(value)

    return None


# ---------------------------------------------------------------------------
# ② 內嵌 workbook
# ---------------------------------------------------------------------------
def read_embedded_series(chart: Any) -> dict[str, list[float | None]]:
    """
    讀出圖表內嵌工作簿的數值，即使用者右鍵「編輯資料」看到的內容。

    python-pptx 寫入的內嵌工作簿佈局：A 欄為類別，B 欄起為各系列，
    第 1 列為系列名稱。散點圖則 A 欄為 x 值、B 欄為 y 值。
    """
    workbook_part = chart.part.chart_workbook.xlsx_part

    if workbook_part is None:
        raise ValueError("圖表沒有內嵌工作簿")

    workbook = load_workbook(io.BytesIO(workbook_part.blob), data_only=True)
    worksheet = workbook.active

    series_values: dict[str, list[float | None]] = {}

    for column_index in range(2, worksheet.max_column + 1):
        header = worksheet.cell(row=1, column=column_index).value
        name = str(header) if header is not None else f"系列{column_index - 1}"

        values = [
            _to_float(worksheet.cell(row=row_index, column=column_index).value)
            for row_index in range(2, worksheet.max_row + 1)
        ]

        series_values[name] = values

    workbook.close()
    return series_values


# ---------------------------------------------------------------------------
# ③ 外部稽核 Excel
# ---------------------------------------------------------------------------
@dataclass
class ExternalSheet:
    """稽核 Excel 中一張工作表的內容。"""

    sheet_name: str
    chart_title: str | None
    metric_key: str | None
    page_number: int | None
    series: dict[str, list[float | None]] = field(default_factory=dict)
    #: 第一欄（類別／標籤）的內容。原生表格靠它做內容比對配對。
    categories: list[str] = field(default_factory=list)
    #: 稽核表記錄的圖表類型（column／table／heatmap…）。
    chart_kind: str | None = None
    #: 稽核表記錄的單位。表格比對要用它套用相同的顯示格式。
    unit: str | None = None


def read_external_workbook(path: str | Path) -> list[ExternalSheet]:
    """
    讀取稽核 Excel。

    表頭位置不固定（前面有指標資訊區），因此：
    - 以第一欄的欄位標籤（「圖表標題」「指標鍵」「簡報頁碼」）讀取 metadata
    - 以第一欄出現「類別」或「標籤」的那一列作為資料表頭列
    """
    workbook = load_workbook(str(path), data_only=True)
    sheets: list[ExternalSheet] = []

    for worksheet in workbook.worksheets:
        if worksheet.title == "索引":
            continue

        header_row = None
        metadata: dict[str, Any] = {}

        for row_index in range(1, min(worksheet.max_row, 40) + 1):
            first = worksheet.cell(row=row_index, column=1).value

            if not isinstance(first, str):
                continue

            label = first.strip()

            if label in _CATEGORY_HEADERS:
                header_row = row_index
                break

            metadata[label] = worksheet.cell(row=row_index, column=2).value

        if header_row is None:
            continue

        page_number = metadata.get("簡報頁碼")

        sheet = ExternalSheet(
            sheet_name=worksheet.title,
            chart_title=(
                str(metadata["圖表標題"]).strip()
                if metadata.get("圖表標題") is not None
                else None
            ),
            metric_key=(
                str(metadata["指標鍵"]).strip()
                if metadata.get("指標鍵") is not None
                else None
            ),
            page_number=(
                int(page_number) if isinstance(page_number, (int, float)) else None
            ),
            chart_kind=(
                str(metadata["圖表類型"]).strip()
                if metadata.get("圖表類型") is not None
                else None
            ),
            unit=(
                str(metadata["單位"]).strip()
                if metadata.get("單位") is not None
                else None
            ),
        )

        sheet.categories = [
            str(value).strip()
            for row_index in range(header_row + 1, worksheet.max_row + 1)
            if (value := worksheet.cell(row=row_index, column=1).value) is not None
        ]

        series_values: dict[str, list[float | None]] = {}

        for column_index in range(2, worksheet.max_column + 1):
            header = worksheet.cell(row=header_row, column=column_index).value

            if header is None:
                continue

            name = str(header).strip()

            # 來源欄不是數值系列，跳過。
            if name.startswith("來源"):
                continue

            values = [
                _to_float(
                    worksheet.cell(row=row_index, column=column_index).value
                )
                for row_index in range(header_row + 1, worksheet.max_row + 1)
            ]

            # 去掉表尾可能的空白列
            while values and values[-1] is None:
                values.pop()

            series_values[name] = values

        sheet.series = series_values
        sheets.append(sheet)

    workbook.close()
    return sheets


def read_table_shape(shape: Any) -> tuple[list[str], list[str], dict[str, list[float | None]]]:
    """
    讀出原生表格的內容。

    表格的數字是**文字**，所以要反向解析回數值才能比對。解析一律用
    :func:`table_builder.parse_value`，與寫入時的 :func:`format_value`
    成對——格式化規則只有一份，比對才有意義。

    Returns:
        (欄名清單, 列標籤清單, {欄名: 數值清單})。
    """
    table = shape.table
    rows = list(table.rows)

    if len(rows) < 2:
        return [], [], {}

    header_cells = list(rows[0].cells)
    column_names = [cell.text.strip() for cell in header_cells[1:]]

    row_labels: list[str] = []
    values: dict[str, list[float | None]] = {name: [] for name in column_names}

    for row in rows[1:]:
        cells = list(row.cells)
        row_labels.append(cells[0].text.strip())

        for index, name in enumerate(column_names, start=1):
            text = cells[index].text if index < len(cells) else ""
            values[name].append(table_builder.parse_value(text))

    return column_names, row_labels, values


def _match_external_sheet_for_table(
    sheets: Sequence[ExternalSheet],
    column_names: Sequence[str],
    row_labels: Sequence[str],
) -> ExternalSheet | None:
    """
    為原生表格找出對應的稽核工作表。

    表格沒有 chart_title 可用（title 在投影片標題上，與稽核表記的
    「圖表標題」不是同一個字串），所以改以內容配對：欄名與列標籤同時
    吻合的工作表就是它。內容配對比序號配對更穩——序號會因封面／章節頁
    偏移而錯頁，曾實際造成誤判。
    """
    wanted_columns = [name for name in column_names]
    wanted_rows = list(row_labels)

    for sheet in sheets:
        if list(sheet.series) != wanted_columns:
            continue

        if sheet.categories[: len(wanted_rows)] == wanted_rows:
            return sheet

    return None


def _match_external_sheet(
    sheets: Sequence[ExternalSheet],
    chart_title: str,
    slide_number: int,
    chart_ordinal: int,
) -> ExternalSheet | None:
    """
    找出與圖表對應的稽核工作表。

    優先使用「實際投影片頁碼 + 圖表標題」這組複合鍵。單用標題並不安全：
    排名頁與原始值頁可能產生相同標題，但引用不同 metric，曾造成實值被拿去
    和排名比較。只有標題在整本稽核 Excel 中唯一時，才允許單獨使用標題。

    舊版產物的頁碼可能是邏輯頁碼而非實際 slide 序號，因此最後仍保留頁碼與
    「第幾張含圖表的投影片」兩層相容性 fallback。
    """
    normalized = chart_title.strip()
    title_matches = [
        sheet
        for sheet in sheets
        if sheet.chart_title and sheet.chart_title.strip() == normalized
    ]

    for sheet in title_matches:
        if sheet.page_number == slide_number:
            return sheet

    if len(title_matches) == 1:
        return title_matches[0]

    page_matches = [
        sheet for sheet in sheets if sheet.page_number == slide_number
    ]

    if len(page_matches) == 1:
        return page_matches[0]

    for sheet in sheets:
        if sheet.page_number == chart_ordinal:
            return sheet

    return None


# ---------------------------------------------------------------------------
# 主驗證流程
# ---------------------------------------------------------------------------
def verify(
    pptx_path: str | Path,
    external_xlsx_path: str | Path | None = None,
) -> VerificationReport:
    """
    執行三方比對。

    Args:
        pptx_path: 待驗證的 .pptx。
        external_xlsx_path: 稽核 Excel。省略時只做①②兩方比對。
    """
    presentation = Presentation(str(pptx_path))
    report = VerificationReport()

    external: list[ExternalSheet] = []

    if external_xlsx_path is not None:
        external = read_external_workbook(external_xlsx_path)

    # pptx 的 slide 序號與稽核 Excel 的頁碼可能有偏移（封面頁），
    # 因此另外維護一個「第幾張含圖表的投影片」計數，作為回退配對依據。
    chart_slide_ordinal = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                _verify_table_shape(shape, slide_number, external, report)
                continue

            if not shape.has_chart:
                continue

            chart = shape.chart
            chart_slide_ordinal += 1

            title = (
                chart.chart_title.text_frame.text
                if chart.has_title
                else f"(無標題圖表 @slide {slide_number})"
            )

            try:
                embedded = read_embedded_series(chart)
            except Exception as error:  # noqa: BLE001
                report.warnings.append(
                    f"slide {slide_number} 圖表 {title!r} 無法讀取內嵌工作簿："
                    f"{error}"
                )
                continue

            external_sheet = None

            if external:
                external_sheet = _match_external_sheet(
                    external, title, slide_number, chart_slide_ordinal
                )

                if external_sheet is None:
                    report.warnings.append(
                        f"slide {slide_number} 圖表 {title!r} "
                        "在稽核 Excel 中找不到對應工作表，僅比對①②"
                    )

            # 走遍**所有** plot，不只 plots[0]。雙軸圖是兩個 plot
            #（c:barChart + c:lineChart），只看第一個會讓折線那一軸的系列
            # 靜默漏驗——而次軸系列正是量級差異大、最需要核對的那組。
            series_ordinal = -1

            for plot in chart.plots:
                for series in plot.series:
                    series_ordinal += 1
                    cache_values = [_to_float(value) for value in series.values]

                    embedded_values = embedded.get(series.name)

                    if embedded_values is None:
                        # 系列名稱對不上時，退回按順序取，避免整筆漏驗。
                        # 序號跨 plot 連續計算：內嵌 workbook 的欄序是
                        # add_chart() 當下的系列順序，與 plot 切分無關。
                        ordered = list(embedded.values())
                        embedded_values = (
                            ordered[series_ordinal]
                            if series_ordinal < len(ordered)
                            else []
                        )
                        report.warnings.append(
                            f"slide {slide_number} 圖表 {title!r} 系列 "
                            f"{series.name!r} 在內嵌工作簿中找不到同名欄，"
                            "已按欄位順序比對"
                        )

                    external_values = None

                    if external_sheet is not None:
                        external_values = _external_values_for(
                            external_sheet, series.name
                        )

                        if external_values is None:
                            report.warnings.append(
                                f"slide {slide_number} 圖表 {title!r} 系列 "
                                f"{series.name!r} 在稽核 Excel 工作表 "
                                f"{external_sheet.sheet_name!r} 中找不到對應欄位"
                            )

                    # 內嵌工作簿的列數可能多於實際資料（python-pptx 不會裁掉
                    # 空白列），比對前先對齊長度。
                    embedded_values = _trim_to(
                        embedded_values, len(cache_values)
                    )

                    if external_values is not None:
                        external_values = _trim_to(
                            external_values, len(cache_values)
                        )

                    report.comparisons.append(
                        SeriesComparison(
                            slide_number=slide_number,
                            chart_title=title,
                            series_name=series.name,
                            chart_cache=cache_values,
                            embedded=embedded_values,
                            external=external_values,
                        )
                    )

    return report


def _verify_table_shape(
    shape: Any,
    slide_number: int,
    external: Sequence[ExternalSheet],
    report: VerificationReport,
) -> None:
    """
    比對原生表格的儲存格文字與稽核 Excel。

    表格沒有內嵌 workbook，因此這裡只有兩份副本可比。少一份副本不等於
    可以不驗——表格上的數字一樣是主管會照著唸的數字。
    """
    column_names, row_labels, values = read_table_shape(shape)

    if not column_names:
        report.warnings.append(
            f"slide {slide_number} 的表格沒有可解析的表頭，未納入比對"
        )
        return

    title = f"(原生表格 @slide {slide_number})"
    sheet = None

    if external:
        sheet = _match_external_sheet_for_table(
            external, column_names, row_labels
        )

        if sheet is None:
            report.warnings.append(
                f"slide {slide_number} 的原生表格在稽核 Excel 中"
                "找不到內容相符的工作表，無法比對"
            )
    else:
        report.warnings.append(
            f"slide {slide_number} 的原生表格沒有內嵌工作簿可比對，"
            "需提供稽核 Excel 才能驗證其數值"
        )

    for name in column_names:
        cell_values = values[name]
        external_values = None

        if sheet is not None:
            external_values = sheet.series.get(name)

            if external_values is not None:
                # 儲存格文字是四捨五入後的顯示值（如 15.6369% → "15.64%"），
                # 直接跟稽核 Excel 的完整精度比一定不等。所以把稽核值套用
                # **同一套** format_value 規則再解析回來，比的才是同一個
                # 數字的同一種寫法。若差異大於顯示精度，這樣仍然抓得到。
                external_values = [
                    table_builder.parse_value(
                        table_builder.format_value(value, sheet.unit),
                        sheet.unit,
                    )
                    for value in _trim_to(external_values, len(cell_values))
                ]

        report.comparisons.append(
            SeriesComparison(
                slide_number=slide_number,
                chart_title=title,
                series_name=name,
                chart_cache=cell_values,
                embedded=None,
                external=external_values,
                is_table=True,
            )
        )


def _external_values_for(
    sheet: ExternalSheet,
    series_name: str,
) -> list[float | None] | None:
    """
    在稽核工作表中找出與圖表系列對應的欄位。

    散點圖的系列名格式為 ``"{x欄} vs {y欄}"``，而稽核表把 x、y 分成兩欄。
    python-pptx 的 ``series.values`` 對散點圖回傳的是 y 值，因此取 y 欄比對。
    """
    if series_name in sheet.series:
        return sheet.series[series_name]

    if " vs " in series_name:
        _, _, y_name = series_name.partition(" vs ")
        y_name = y_name.strip()

        if y_name in sheet.series:
            return sheet.series[y_name]

    # 單欄工作表沒有歧義，可安全地按唯一欄位配對。
    if len(sheet.series) == 1:
        return next(iter(sheet.series.values()))

    return None


def _trim_to(values: Sequence[float | None], length: int) -> list[float | None]:
    """截斷或補 None 到指定長度，便於等長比對。"""
    trimmed = list(values[:length])

    while len(trimmed) < length:
        trimmed.append(None)

    return trimmed


def print_report(report: VerificationReport) -> None:
    """輸出人類可讀的驗證報告。"""
    print("=" * 72)
    print("三方數值比對驗證（T1）")
    print("=" * 72)

    for comparison in report.comparisons:
        status = "PASS" if comparison.passed else "FAIL"
        if comparison.is_table:
            external_note = (
                "原生表格：儲存格文字 ↔ 稽核 Excel"
                if comparison.external is not None
                else "原生表格：無可比對副本"
            )
        else:
            external_note = (
                "含外部稽核 Excel"
                if comparison.external is not None
                else "未比對外部 Excel"
            )

        print(
            f"[{status}] slide {comparison.slide_number} "
            f"{comparison.chart_title!r} / 系列 {comparison.series_name!r} "
            f"（{external_note}）"
        )

        if not comparison.passed:
            print(f"        {comparison.describe_failure()}")

    if report.warnings:
        print("\n警告：")

        for warning in report.warnings:
            print(f"  - {warning}")

    total = len(report.comparisons)
    failed = len(report.failures)

    print("\n" + "-" * 72)
    print(
        f"比對系列數：{total}｜通過：{total - failed}｜失敗：{failed}｜"
        f"含外部 Excel 比對：{report.external_checked}"
    )
    print(
        "最終結果：",
        "全部一致 (PASS)" if report.passed else "發現不一致 (FAIL)",
    )
    print("-" * 72)


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="驗證 PPT 圖表的三份資料副本數值是否一致"
    )
    parser.add_argument("pptx", help="待驗證的 .pptx 路徑")
    parser.add_argument(
        "xlsx",
        nargs="?",
        default=None,
        help="稽核用 .xlsx 路徑（省略則只比對圖表快取與內嵌工作簿）",
    )

    args = parser.parse_args(argv)

    report = verify(args.pptx, args.xlsx)
    print_report(report)

    return 0 if report.passed else 1


if __name__ == "__main__":
    sys.exit(main())

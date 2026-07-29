"""
圖表生成模組 (POC)
=====================
核心設計原則：
1. 圖表數字的「唯一真相來源」是本模組的輸入資料（解析後的 Excel 結構化資料）。
2. 一律透過 python-pptx 的 add_chart() API 生成圖表，
   讓「畫面顯示值 (chart cache)」與「右鍵編輯資料看到的內嵌工作表」
   自動保持一致，不手動操作底層 XML。
3. 不使用外部連結 (external link) 方式綁定 Excel，
   因為 PPT 寄出後外部路徑會失效。改用 PowerPoint 原生的
   「內嵌工作簿 (embedded workbook)」機制，這正是使用者
   右鍵「編輯資料」時開啟的資料來源。
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Emu


@dataclass
class ChartSpec:
    """一張圖表所需的資料規格。"""

    title: str
    categories: Sequence[str]
    series: dict[str, Sequence[float]]  # {系列名稱: 數值}
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED


def add_category_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    將 ChartSpec 的資料寫入圖表，並插入到指定投影片。

    重點：CategoryChartData 是唯一資料入口。
    python-pptx 會同時：
      1. 寫入 chart1.xml 的 <c:numCache>（畫面顯示用）
      2. 產生對應的 embedded .xlsx（右鍵「編輯資料」開啟用）
    兩者保證數值一致，因為都是從這裡的同一份 spec 產生。
    """
    chart_data = CategoryChartData()
    chart_data.categories = list(spec.categories)
    for series_name, values in spec.series.items():
        chart_data.add_series(series_name, values)

    graphic_frame = slide.shapes.add_chart(
        spec.chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = spec.title
    return chart


@dataclass
class ScatterSpec:
    """散點圖（規模 vs 成長）專用規格。"""

    title: str
    series_name: str
    points: Sequence[tuple[float, float]]  # [(x, y), ...]
    labels: Sequence[str] | None = None  # 對應每個點的標籤（如銀行名稱）


def add_scatter_chart(
    slide: Slide,
    spec: ScatterSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    散點圖走 XyChartData，一樣是唯一資料入口原則。
    labels 部分 python-pptx 原生不支援資料點標籤文字，
    若模板要求顯示銀行名稱標籤，需要在 chart 生成後
    透過底層 XML（c:dLbls）額外處理，這裡先留 TODO。
    """
    chart_data = XyChartData()
    series = chart_data.add_series(spec.series_name)
    for x, y in spec.points:
        series.add_data_point(x, y)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = spec.title
    # TODO: 若需散點標籤（如銀行名稱），需手動操作 chart.plots[0] 的
    # dLbls XML 節點，python-pptx 目前無高階 API。
    return chart


def add_pie_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(6096000),
    height: int = Emu(4114800),
):
    """
    市占率圖（圓餅圖）。走同一個 CategoryChartData 入口，
    僅 chart_type 固定為 PIE。要求 spec.series 只能有一組系列。
    """
    if len(spec.series) != 1:
        raise ValueError("圓餅圖只能有一組系列，請確認 ChartSpec.series 長度為 1")
    pie_spec = ChartSpec(
        title=spec.title,
        categories=spec.categories,
        series=spec.series,
        chart_type=XL_CHART_TYPE.PIE,
    )
    return add_category_chart(slide, pie_spec, left, top, width, height)


# ---------------------------------------------------------------------------
# Skill Registry：圖表 Agent 只透過 skill 名稱字串呼叫，不需知道實作細節。
# 新增圖表類型時，在此註冊即可，Agent 端無需改動。
# ---------------------------------------------------------------------------
CHART_SKILLS = {
    "column": add_category_chart,   # 長條圖／排名圖／成長率圖
    "bar": add_category_chart,      # 橫條圖（呼叫時 spec.chart_type 需先設為 BAR_CLUSTERED）
    "pie": add_pie_chart,           # 市占率圖
    "scatter": add_scatter_chart,   # 規模 vs 成長散點圖
    # "heatmap": 由原生表格 + 儲存格底色模擬，不走 add_chart()，
    #            另外註冊在 table_builder.py（尚未實作），Agent 需個別處理。
}


# ---------------------------------------------------------------------------
# Function Calling / Tool Use 對應層
# ---------------------------------------------------------------------------
# LLM API（OpenAI / Anthropic / Bedrock 等）本身無法執行 Python 檔案，
# 它只能根據我們提供的 "tool schema" 回傳一段結構化 JSON，
# 表達「我想呼叫哪個 skill、帶哪些參數」。
#
# 真正的執行流程是：
#   1. 我們把 CHART_SKILLS 的 key 轉成 tool schema，隨 prompt 送給 LLM
#   2. LLM 回傳 {"name": "pie", "arguments": {...}}（純文字/JSON，不是程式碼）
#   3. 我們自己的 dispatcher（如下 dispatch_chart_skill）收到這段 JSON，
#      在本地 Python 環境查表、驗證、實際呼叫對應函式
#
# LLM 從未執行過任何一行本地程式碼，執行權限完全留在我方。
# ---------------------------------------------------------------------------

CHART_SKILL_TOOL_SCHEMAS = [
    {
        "name": "column",
        "description": "長條圖，適合排名、成長率比較等類別型數據。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {
                    "type": "string",
                    "description": "MetricStore 中的指標鍵，不可填入實際數值。",
                },
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "限定取用的系列名稱，留空代表全取。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "chart_title"],
        },
    },
    {
        "name": "pie",
        "description": "圓餅圖，適合市占率等單一系列的佔比數據。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {"type": "string"},
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "chart_title"],
        },
    },
    {
        "name": "scatter",
        "description": "散點圖，適合「規模 vs 成長」等雙變數關係。",
        "parameters": {
            "type": "object",
            "properties": {
                "metric_key": {"type": "string"},
                "series_names": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "恰好兩個系列名稱，分別對應 x 軸與 y 軸。",
                },
                "chart_title": {"type": "string"},
            },
            "required": ["metric_key", "series_names", "chart_title"],
        },
    },
]
"""
給 LLM API 的 tool schema 清單。注意 schema 中完全沒有 "values" 這種
數字欄位，LLM 只能填 metric_key 引用，這是刻意設計，逼迫 LLM 無法
自己編造數字（呼應 .kiro/steering/tech.md 第 1 條原則）。

實際串接時（OpenAI 範例）：

    response = client.chat.completions.create(
        model=...,
        messages=[...],
        tools=[{"type": "function", "function": schema}
               for schema in CHART_SKILL_TOOL_SCHEMAS],
    )
    tool_call = response.choices[0].message.tool_calls[0]
    skill_name = tool_call.function.name          # 例如 "pie"
    raw_args = json.loads(tool_call.function.arguments)

Anthropic / Bedrock 的 tool use 格式略有不同（key 名稱不同），
但核心流程一致：LLM 回傳「工具名稱 + 參數」，執行仍在本地端。
"""


def dispatch_chart_skill(skill_name: str) -> callable:
    """
    Dispatcher：LLM 回傳的 skill_name（純字串）透過這裡查表，
    取得實際的 Python 函式並在本地執行。

    這是 LLM 輸出與本地程式碼執行之間唯一的橋接點，
    也是防呆的關卡：LLM 若回傳未註冊的名稱，直接拋錯，
    不會有機會執行任何非預期的程式碼路徑
    （因為只從 CHART_SKILLS 這個白名單裡取值，不會 eval/exec LLM 的輸出）。
    """
    skill = CHART_SKILLS.get(skill_name)
    if skill is None:
        raise ValueError(
            f"LLM 回傳了未註冊的 skill: {skill_name!r}，"
            f"可用選項: {list(CHART_SKILLS.keys())}"
        )
    return skill

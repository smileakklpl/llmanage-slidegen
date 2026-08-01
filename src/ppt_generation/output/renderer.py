"""
PPT 組裝
=========
對應 docs/圖表原生性與資料同步設計.md Stage 5 與 §9。

職責：套用 `source/template.pptx`，逐頁組裝

1. 標題（來自 SectionPlan）
2. 原生圖表（一律經 ``CHART_SKILLS`` → ``add_chart()`` 單一入口）
3. 敘事文字（佔位符在此代入實際數值）

版面策略：模板的 `1_標題及內容` 版面只有一個全寬內容區，本模組把它
縮成右側欄位放文字，左側留給圖表，形成顧問簡報常見的「圖左文右」佈局。

**不可退讓的原則**：本模組不得自行寫入任何數字。圖表數值來自
ChartSpec（由 MetricStore 查表產生），敘事數值來自
:func:`placeholders.render_text` 查表代入。
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Sequence

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..core import config, placeholders, theme
from ..agents.narrative_writer import PageNarrative
from ..agents.section_planner import CONCLUSION_CHAPTER, SectionPlan
from ..charts.chart_builder import ScatterSpec
from ..charts.chart_planner import (
    TABLE_LIKE_CHARTS,
    VISUAL_SKILLS,
    ChartPlan,
    ResolvedChart,
    resolve_chart_plan,
)
from ..contracts import DeckSpecContract
from ..data.metric_store import MetricStore


logger = logging.getLogger(__name__)

#: 模板中用於內容頁的版面名稱。
CONTENT_LAYOUT_NAME = "1_標題及內容"

#: 模板中用於章節分隔頁的版面名稱。
SECTION_LAYOUT_NAME = "2_章節標題"

#: 模板中用於結尾頁的版面名稱（附件一模板第 5 頁用的就是這個）。
CLOSING_LAYOUT_NAME = "3_標題投影片"

#: 目錄頁標題與結尾頁文字。都可由呼叫端覆寫。
AGENDA_TITLE = "目錄"
CLOSING_MESSAGE = "感謝聆聽"

#: 封面頁在模板中是第 1 張投影片，保留不刪；其後的示範頁移除。
#: 「封面 + 目錄」共 2 張非內容頁，是頁碼推算的固定前綴。
FRONT_MATTER_SLIDES = 2

#: placeholder 索引（由模板結構決定，見 template.pptx 版面配置）。
PH_TITLE = 0
PH_BODY = 1

#: 圖表占內容區寬度的比例，其餘留給敘事文字。
CHART_WIDTH_RATIO = 0.60


class RenderError(RuntimeError):
    """PPT 組裝失敗。"""


@dataclass
class PageBundle:
    """一頁所需的全部素材。"""

    section: SectionPlan
    chart: ResolvedChart | None = None
    narrative: PageNarrative | None = None


@dataclass
class RenderReport:
    """組裝結果摘要。"""

    output_path: Path | None = None
    page_count: int = 0
    chart_count: int = 0
    #: 佔位符代入失敗的訊息（頁碼 → 錯誤清單）
    placeholder_errors: dict[int, list[str]] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)
    #: 章節分隔頁數量
    divider_count: int = 0
    #: 最終 .pptx 的實際投影片總數（含封面、目錄、章節頁、結尾頁）
    slide_count: int = 0
    #: 依序的章節名稱
    chapters: list[str] = field(default_factory=list)
    #: 是否產出了結論頁
    conclusion_page: bool = False


# ---------------------------------------------------------------------------
# 頁碼指派
# ---------------------------------------------------------------------------
def chapter_order(bundles_or_sections: Sequence[Any]) -> list[str]:
    """依出現順序取出章節名稱（跳過未歸章節者）。"""
    chapters: list[str] = []

    for item in bundles_or_sections:
        section = item.section if isinstance(item, PageBundle) else item

        if section.chapter and section.chapter not in chapters:
            chapters.append(section.chapter)

    return chapters


def assign_page_numbers(
    sections: Sequence[SectionPlan],
    *,
    include_agenda: bool = True,
) -> list[SectionPlan]:
    """
    把每個內容頁的 ``page_number`` 改成它在最終 .pptx 中的實際投影片序號。

    這件事必須在圖表決策**之前**做完：稽核 Excel 的工作表名是
    ``P.{頁碼}_{指標名稱}``（FR-3.1），而頁碼是從 ChartPlan 帶下去的。
    若這裡的頁碼是「第幾個內容頁」而不是「第幾張投影片」，主管拿著
    Excel 對照簡報就會翻錯頁——封面、目錄與每張章節分隔頁都會造成偏移。

    版面順序：封面 → 目錄 →（章節分隔頁 → 該章節內容頁…）× N → 結尾頁。

    就地修改並回傳同一批物件，方便鏈式使用。
    """
    # 沒有任何頁面歸屬章節時 render_deck 不會產目錄頁（目錄會是空的），
    # 這裡的推算必須跟著它，否則頁碼會整份差一頁。
    has_chapters = any(section.chapter for section in sections)

    number = 1 + (1 if (include_agenda and has_chapters) else 0)  # 封面（+ 目錄）
    current_chapter: str | None = None

    for section in sections:
        if section.chapter and section.chapter != current_chapter:
            number += 1  # 章節分隔頁
            current_chapter = section.chapter

        number += 1
        section.page_number = number

    return list(sections)


# ---------------------------------------------------------------------------
# 版面計算
# ---------------------------------------------------------------------------
@dataclass(frozen=True)
class ContentArea:
    """內容區座標（EMU）。由模板的 BODY placeholder 推導。"""

    left: int
    top: int
    width: int
    height: int

    def split(
        self,
        chart_ratio: float = CHART_WIDTH_RATIO,
        gutter: int = Inches(0.25),
    ) -> tuple[ContentArea, ContentArea]:
        """切成左（圖表）右（文字）兩欄。"""
        chart_width = int((self.width - gutter) * chart_ratio)
        text_width = self.width - gutter - chart_width

        chart_area = ContentArea(self.left, self.top, chart_width, self.height)
        text_area = ContentArea(
            self.left + chart_width + gutter,
            self.top,
            text_width,
            self.height,
        )

        return chart_area, text_area


def _content_area(layout: Any) -> ContentArea:
    """從版面的 BODY placeholder 取得內容區範圍。"""
    for placeholder in layout.placeholders:
        if placeholder.placeholder_format.idx == PH_BODY:
            return ContentArea(
                left=placeholder.left,
                top=placeholder.top,
                width=placeholder.width,
                height=placeholder.height,
            )

    raise RenderError(
        f"版面 {layout.name!r} 找不到 BODY placeholder（idx={PH_BODY}），"
        "無法推導內容區範圍"
    )


def _find_layout(presentation: Presentation, name: str) -> Any:
    for layout in presentation.slide_layouts:
        if layout.name == name:
            return layout

    available = [layout.name for layout in presentation.slide_layouts]
    raise RenderError(f"模板中找不到版面 {name!r}。可用版面：{available}")


# ---------------------------------------------------------------------------
# 文字填充
# ---------------------------------------------------------------------------
def _write_segments(
    paragraph: Any,
    segments: Sequence[placeholders.TextSegment],
    *,
    size: Any,
    color: Any,
    bold: bool = False,
) -> None:
    """
    把代入後的片段寫成一連串 run，來自 MetricStore 的值加黃色標示。

    一段文字為什麼要拆成多個 run：黃色標示是 run 級屬性。整段套用會把
    敘事文字也一起標黃，那就不是「重點」而是塗滿；只標值的話，讀者一眼
    看到的黃色字元恰好等於系統算出來的數字。
    """
    for segment in segments:
        if not segment.text:
            continue

        run = paragraph.add_run()
        run.text = segment.text

        theme.apply_font(
            run,
            size=size,
            bold=bold or segment.from_metric,
            color=theme.TITLE_COLOR if segment.from_metric else color,
            highlight=theme.HIGHLIGHT if segment.from_metric else None,
        )


def _fill_narrative(
    placeholder: Any,
    narrative: PageNarrative,
    store: MetricStore,
    *,
    include_headline: bool = False,
) -> list[str]:
    """
    把敘事要點填入文字框，佔位符在此代入實際數值。

    ``include_headline=False``：headline 已由重點訊息帶呈現（附件三的版面
    語彙），這裡只放要點，避免同一句話在一頁裡出現兩次。

    ``strict=False``：單一佔位符失敗不應讓整份簡報生不出來，
    改為保留原佔位符並回報錯誤，讓使用者一眼看出哪裡沒接上。
    """
    errors: list[str] = []
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    # 「溢出時縮小文字」的保險。字數上限由 narrative_writer 守著，這裡是
    # 第二道防線：真實模型偶爾會寫出貼著上限的長句，寧可字小一點，
    # 也不要讓文字流到投影片外面被裁掉。
    #
    # 但它只會縮小、不會放大——短敘事仍停在 11pt，在 40% 寬的文字欄裡
    # 看起來像忘了排版。所以字級改為先算出來，autofit 只留作最後保險。
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    lines: list[tuple[str, bool]] = []

    if include_headline:
        lines.append((narrative.headline, True))

    lines.extend((bullet, False) for bullet in narrative.bullets)

    bullet_size, headline_size = _narrative_font_sizes(placeholder, lines)

    first = True

    for text, is_headline in lines:
        segments, segment_errors = placeholders.render_segments(
            text, store, strict=False
        )
        errors.extend(segment_errors)

        paragraph = (
            text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        )
        first = False
        paragraph.level = 0

        if not is_headline:
            # 項目符號自己寫，不依賴版面的 buChar：模板的內容版面在不同
            # 縮排層級會給出不同符號，一頁裡混用兩種符號很難看。
            marker = paragraph.add_run()
            marker.text = theme.ACTION_BULLET_PREFIX
            theme.apply_font(
                marker,
                size=bullet_size,
                bold=True,
                color=theme.ACCENT,
            )

        _write_segments(
            paragraph,
            segments,
            size=headline_size if is_headline else bullet_size,
            color=theme.TITLE_COLOR if is_headline else theme.BODY_COLOR,
            bold=is_headline,
        )

        # 段距隨字級縮放。字級縮到 9pt 卻還留 8pt 段距，等於把省下來的
        # 高度又花在空白上。比例與 theme 估算高度時用的同一個常數，
        # 兩邊不一致的話估算會愈偏愈多。
        paragraph.space_after = theme.paragraph_gap_for(bullet_size)

    return errors


def _narrative_font_sizes(
    placeholder: Any,
    lines: Sequence[tuple[str, bool]],
) -> tuple[Any, Any]:
    """
    依敘事總長度與文字框大小推算要點與 headline 的字級。

    兩者一起算而不是各算一次：它們共用同一個文字框的高度，分開算會各自
    以為自己有整框可用，加起來就溢出了。做法是先算出要點字級，headline
    再按原本的比例（13:11）往上抬，並夾在上限內。

    Returns:
        ``(要點字級, headline 字級)``。
    """
    # 項目符號也占寬度，估算時一併計入。
    combined = "".join(
        text if is_headline else f"{theme.ACTION_BULLET_PREFIX}{text}"
        for text, is_headline in lines
    )

    bullet_size = theme.fit_font_size(
        combined,
        placeholder.width,
        placeholder.height,
        maximum=theme.MAX_BODY_FONT_SIZE,
        minimum=theme.MIN_BODY_FONT_SIZE,
        paragraph_count=len(lines),
    )

    ratio = int(theme.HEADLINE_FONT_SIZE) / int(theme.BULLET_FONT_SIZE)
    headline_size = Pt(
        min(
            int(bullet_size) * ratio / 12700,
            int(theme.HEADLINE_FONT_SIZE) / 12700,
        )
    )

    return bullet_size, headline_size


def _style_title(slide: Any, text: str) -> None:
    """
    標題一律左上、深灰，與附件三一致；字級依標題長度推算。

    22pt 只是上限。真實模型寫出的頁標題長度差距很大（「市場概況」四個字
    到「有效卡率與循環信用餘額的交叉觀察」十六個字），固定 22pt 的結果是
    短標題看起來空、長標題折成兩行把重點訊息帶往下推。
    """
    title_shape = slide.shapes.title

    if title_shape is None:
        return

    text_frame = title_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT

    run = paragraph.add_run()
    run.text = text

    theme.apply_font(
        run,
        size=theme.fit_single_line_font_size(
            text,
            title_shape.width,
            maximum=theme.TITLE_FONT_SIZE,
            minimum=theme.MIN_TITLE_FONT_SIZE,
        ),
        bold=False,
        color=theme.TITLE_COLOR,
    )


def add_key_message_bar(
    slide: Any,
    text: str,
    store: MetricStore,
    area: ContentArea,
) -> tuple[Any, list[str]]:
    """
    在內容區頂端加一條重點訊息帶，放本頁的結論句。

    這是附件三每一頁都有的元件（「◆ 市場總流通卡數穩定成長至…」），
    也是顧問簡報「一頁一結論」的載體：讀者只看這一行就該知道本頁結論。
    底色為淡黃、左緣一條純黃標示條，句中的數值再加黃色螢光標示。

    Returns:
        (訊息帶 shape, 佔位符錯誤清單)。
    """
    marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(area.left),
        Emu(area.top),
        theme.KEY_BAR_MARKER_WIDTH,
        theme.KEY_BAR_HEIGHT,
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = theme.KEY_BAR_MARKER
    marker.line.fill.background()
    marker.shadow.inherit = False

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(area.left + theme.KEY_BAR_MARKER_WIDTH),
        Emu(area.top),
        Emu(area.width - theme.KEY_BAR_MARKER_WIDTH),
        theme.KEY_BAR_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.KEY_BAR_FILL
    bar.line.fill.background()
    bar.shadow.inherit = False

    text_frame = bar.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = theme.KEY_BAR_TEXT_INSET
    text_frame.margin_right = Emu(91440)
    text_frame.margin_top = Emu(45720)
    text_frame.margin_bottom = Emu(45720)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT

    segments, errors = placeholders.render_segments(text, store, strict=False)

    # 訊息帶是固定高度的單行元件，字級必須依句長算——這是全頁最容易溢出
    # 的地方。估算長度用「代入後」的文字：佔位符 {{流通卡數.value|…|max}}
    # 有二十多個字元，代入後可能只剩「4,512,345」九個，用原字串會把字級
    # 壓到不必要的小。
    resolved_length = theme.KEY_MESSAGE_PREFIX + "".join(
        segment.text for segment in segments
    )

    size = theme.fit_single_line_font_size(
        resolved_length,
        int(area.width) - int(theme.KEY_BAR_MARKER_WIDTH),
        maximum=theme.KEY_MESSAGE_FONT_SIZE,
        minimum=theme.MIN_KEY_MESSAGE_FONT_SIZE,
        inset_x=int(theme.KEY_BAR_TEXT_INSET) + 91440,
    )

    prefix = paragraph.add_run()
    prefix.text = theme.KEY_MESSAGE_PREFIX
    theme.apply_font(
        prefix,
        size=size,
        bold=True,
        color=theme.ACCENT,
    )

    _write_segments(
        paragraph,
        segments,
        size=size,
        color=theme.TITLE_COLOR,
        bold=True,
    )

    return bar, errors


def add_footnote(slide: Any, text: str, area: ContentArea) -> Any:
    """
    在頁尾加一行小字註記。

    圖表頁用它寫出驗收動作（右鍵編輯資料），表格頁改指向稽核 Excel——
    表格沒有內嵌工作簿，沿用同一句話等於在簡報上寫下不成立的承諾。
    """
    box = slide.shapes.add_textbox(
        Emu(area.left),
        theme.FOOTNOTE_TOP,
        Emu(area.width),
        theme.FOOTNOTE_HEIGHT,
    )
    text_frame = box.text_frame
    # 之前是 word_wrap=False，長註記會直接跑出投影片右緣。註記寧可縮小
    # 也不該跑出版面，所以改為允許折行並依長度收字級。
    text_frame.word_wrap = True

    run = text_frame.paragraphs[0].add_run()
    run.text = text

    theme.apply_font(
        run,
        size=theme.fit_single_line_font_size(
            text,
            int(area.width),
            maximum=theme.FOOTNOTE_FONT_SIZE,
            minimum=Pt(7),
        ),
        bold=False,
        color=theme.MUTED_COLOR,
    )

    return box


def _place_text_area(placeholder: Any, area: ContentArea) -> None:
    """把 BODY placeholder 移到指定欄位。"""
    placeholder.left = area.left
    placeholder.top = area.top
    placeholder.width = area.width
    placeholder.height = area.height


# ---------------------------------------------------------------------------
# 頁面組裝
# ---------------------------------------------------------------------------
def add_section_divider(
    presentation: Presentation,
    title: str,
    *,
    index: int | None = None,
) -> Any:
    """
    新增章節分隔頁：``CHAPTER 0N``（紅）+ 章節名 + 紅色底線。

    ``index`` 是章節序號（1 起算）。附件三的章節頁靠這個編號讓聽眾知道
    「講到第幾段了」，沒有編號的章節頁只是一張大字報。
    """
    layout = _find_layout(presentation, SECTION_LAYOUT_NAME)
    slide = presentation.slides.add_slide(layout)
    title_shape = slide.shapes.title

    if title_shape is None:
        return slide

    text_frame = title_shape.text_frame
    text_frame.clear()

    if index is not None:
        label_paragraph = text_frame.paragraphs[0]
        label_run = label_paragraph.add_run()
        label_run.text = f"CHAPTER {index:02d}"
        theme.apply_font(
            label_run,
            size=theme.CHAPTER_LABEL_FONT_SIZE,
            bold=False,
            color=theme.ACCENT,
        )
        title_paragraph = text_frame.add_paragraph()
    else:
        title_paragraph = text_frame.paragraphs[0]

    title_run = title_paragraph.add_run()
    title_run.text = title
    theme.apply_font(
        title_run,
        size=theme.CHAPTER_TITLE_FONT_SIZE,
        bold=True,
        color=theme.BODY_COLOR,
    )

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(title_shape.left),
        Emu(title_shape.top + title_shape.height),
        theme.CHAPTER_RULE_WIDTH,
        theme.CHAPTER_RULE_HEIGHT,
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False

    return slide


def add_agenda_page(
    presentation: Presentation,
    chapters: Sequence[str],
    *,
    title: str = AGENDA_TITLE,
) -> Any:
    """
    新增目錄頁，逐條列出章節。

    目錄的內容完全來自章節清單，不另外撰寫——目錄與章節分隔頁不一致
    是簡報最容易出現、也最傷信任的瑕疵。
    """
    layout = _find_layout(presentation, CONTENT_LAYOUT_NAME)
    slide = presentation.slides.add_slide(layout)
    _style_title(slide, title)

    body = _body_placeholder(slide)

    if body is None:
        return slide

    text_frame = body.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    # 算出來的字級是主要手段，autofit 只是最後保險（章節名意外很長時）。
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # 章節數是變數（本專案 2～8 章都出現過），字級因此不能是常數：
    # 八章節排在 16pt 會頂到頁尾，兩章節排在 16pt 又只占上面四分之一。
    size = theme.fit_font_size(
        "".join(f"00　{chapter}" for chapter in chapters),
        body.width,
        body.height,
        maximum=theme.AGENDA_FONT_SIZE,
        minimum=theme.MIN_BODY_FONT_SIZE,
        paragraph_count=len(chapters),
    )

    for index, chapter in enumerate(chapters):
        paragraph = (
            text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        )
        paragraph.level = 0
        paragraph.space_after = theme.paragraph_gap_for(size)

        number = paragraph.add_run()
        number.text = f"{index + 1:02d}　"
        theme.apply_font(
            number,
            size=size,
            bold=True,
            color=theme.ACCENT,
        )

        label = paragraph.add_run()
        label.text = chapter
        theme.apply_font(
            label,
            size=size,
            bold=False,
            color=theme.TITLE_COLOR,
        )

    return slide


def add_conclusion_page(
    presentation: Presentation,
    bundles: Sequence[PageBundle],
    store: MetricStore,
    *,
    title: str = CONCLUSION_CHAPTER,
) -> tuple[Any, list[str]]:
    """
    新增結論頁：把各章節的結論句收攏成一頁。

    做法刻意是**確定性**的——每個章節取其首頁的 headline，不另外呼叫 LLM
    重寫。理由有兩個：多一次 LLM 呼叫就多一次「結論與內文不一致」的機會；
    而且結論頁的每一句都能指回它出自哪一頁，主管問「這句從哪來的」時
    答案就在簡報裡。數值仍由佔位符代入，黃色標示照樣成立。

    Returns:
        (slide, 佔位符錯誤清單)。
    """
    layout = _find_layout(presentation, CONTENT_LAYOUT_NAME)
    slide = presentation.slides.add_slide(layout)
    _style_title(slide, title)

    errors: list[str] = []
    body = _body_placeholder(slide)

    if body is None:
        return slide, errors

    seen: set[str] = set()
    picked: list[tuple[str, PageNarrative]] = []

    for bundle in bundles:
        chapter = bundle.section.chapter

        if not chapter or chapter in seen or bundle.narrative is None:
            continue

        seen.add(chapter)
        picked.append((chapter, bundle.narrative))

    area = _content_area(layout)
    _place_text_area(body, area)

    text_frame = body.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if not picked:
        return slide, errors

    # 這頁是全份簡報最擁擠的一頁：每個章節一句結論，八章節就是八句長句
    # 全部塞進一個文字框。固定字級在章節多時必定溢出。
    resolved: list[tuple[str, list[placeholders.TextSegment], str]] = []

    for chapter, narrative in picked:
        segments, segment_errors = placeholders.render_segments(
            narrative.headline, store, strict=False
        )
        errors.extend(segment_errors)
        resolved.append(
            (
                f"{theme.ACTION_BULLET_PREFIX}{chapter}｜",
                segments,
                f"（P.{narrative.page_number}）",
            )
        )

    size = theme.fit_font_size(
        "".join(
            label + "".join(s.text for s in segments) + source
            for label, segments, source in resolved
        ),
        body.width,
        body.height,
        maximum=theme.BULLET_FONT_SIZE,
        minimum=theme.MIN_BODY_FONT_SIZE,
        paragraph_count=len(resolved),
    )

    for index, (label_text, segments, source_text) in enumerate(resolved):
        paragraph = (
            text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        )
        paragraph.level = 0
        paragraph.space_after = theme.paragraph_gap_for(size)

        label = paragraph.add_run()
        label.text = label_text
        theme.apply_font(
            label,
            size=size,
            bold=True,
            color=theme.ACCENT,
        )

        _write_segments(
            paragraph,
            segments,
            size=size,
            color=theme.BODY_COLOR,
        )

        source = paragraph.add_run()
        source.text = source_text
        theme.apply_font(
            source,
            size=size,
            bold=False,
            color=theme.MUTED_COLOR,
        )

    return slide, errors


def add_closing_page(
    presentation: Presentation,
    message: str = CLOSING_MESSAGE,
) -> Any:
    """新增結尾頁（模板第 5 頁用的版面）。"""
    layout = _find_layout(presentation, CLOSING_LAYOUT_NAME)
    slide = presentation.slides.add_slide(layout)

    if slide.shapes.title is not None:
        text_frame = slide.shapes.title.text_frame
        text_frame.clear()
        run = text_frame.paragraphs[0].add_run()
        run.text = message
        theme.apply_font(
            run,
            size=theme.CHAPTER_TITLE_FONT_SIZE,
            bold=True,
            color=theme.BODY_COLOR,
        )

    return slide


def add_content_page(
    presentation: Presentation,
    bundle: PageBundle,
    store: MetricStore,
) -> tuple[Any, list[str]]:
    """
    新增一頁內容頁：標題 + 原生圖表 + 敘事。

    Returns:
        (slide, 佔位符錯誤清單)。
    """
    layout = _find_layout(presentation, CONTENT_LAYOUT_NAME)
    slide = presentation.slides.add_slide(layout)
    errors: list[str] = []

    title_text = bundle.section.title or (
        bundle.chart.plan.chart_title if bundle.chart else ""
    )

    _style_title(slide, title_text)

    full_area = _content_area(layout)
    body = _body_placeholder(slide)

    # 重點訊息帶吃掉內容區頂端一條，其餘才是圖表與要點的空間。
    area = full_area

    if bundle.narrative is not None and bundle.narrative.headline.strip():
        _, bar_errors = add_key_message_bar(
            slide, bundle.narrative.headline, store, full_area
        )
        errors.extend(bar_errors)

        offset = int(theme.KEY_BAR_HEIGHT) + int(theme.KEY_BAR_GAP)
        area = ContentArea(
            left=full_area.left,
            top=full_area.top + offset,
            width=full_area.width,
            height=full_area.height - offset,
        )

    if bundle.chart is None:
        # 無圖表的純文字頁，要點占滿剩餘內容區。
        if body is not None and bundle.narrative is not None:
            _place_text_area(body, area)
            errors.extend(_fill_narrative(body, bundle.narrative, store))

        return slide, errors

    chart_area, text_area = area.split()

    if body is not None:
        if bundle.narrative is None:
            # 沒有敘事就移除空的文字框，避免留下空白版面提示文字。
            body._element.getparent().remove(body._element)
        else:
            _place_text_area(body, text_area)
            errors.extend(_fill_narrative(body, bundle.narrative, store))
    elif bundle.narrative is not None:
        errors.append(
            f"第 {bundle.section.page_number} 頁找不到 BODY placeholder，"
            "敘事文字未能填入"
        )

    _insert_chart(slide, bundle.chart, chart_area)

    add_footnote(
        slide,
        (
            theme.TABLE_FOOTNOTE
            if bundle.chart.skill_name in TABLE_LIKE_CHARTS
            else theme.CHART_FOOTNOTE
        ),
        full_area,
    )

    return slide, errors


def _body_placeholder(slide: Any) -> Any | None:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == PH_BODY:
            return shape

    return None


def _insert_chart(
    slide: Any,
    chart: ResolvedChart,
    area: ContentArea,
) -> Any:
    """
    透過 registry 插入原生圖表或原生表格。

    圖表走 ``add_chart()``——這是全系統唯一的圖表落地路徑，會同時寫入
    chart XML 快取與內嵌 workbook，兩者天生一致。表格走 ``add_table()``，
    產出的是真正的 PPT table 物件（FR-2.4：禁止文字方塊拼貼）。
    """
    skill = VISUAL_SKILLS.get(chart.skill_name)

    if skill is None:
        raise RenderError(
            f"圖表 skill {chart.skill_name!r} 未註冊，"
            f"可用選項：{sorted(VISUAL_SKILLS)}"
        )

    return skill(
        slide,
        chart.spec,
        left=Emu(area.left),
        top=Emu(area.top),
        width=Emu(area.width),
        height=Emu(area.height),
    )


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------
def render_deck(
    bundles: Sequence[PageBundle],
    store: MetricStore,
    *,
    output_path: str | Path | None = None,
    template_path: str | Path | None = None,
    deck_title: str | None = None,
    keep_template_slides: bool = False,
    include_agenda: bool = True,
    include_conclusion: bool = True,
    include_closing: bool = True,
    closing_message: str = CLOSING_MESSAGE,
) -> RenderReport:
    """
    組裝整份簡報。

    Args:
        bundles: 每頁的素材。
        store: 唯一真相來源，用於代入敘事佔位符。
        output_path: 輸出 .pptx 路徑，預設 ``outputs/deck.pptx``。
        template_path: 模板路徑，預設 ``source/template.pptx``。
        deck_title: 有值時覆寫模板首頁（封面）標題。
        keep_template_slides: 是否保留模板原有的示範頁。預設 False，
            但模板頁的刪除需操作底層 XML，故目前僅保留首頁並記錄警告。
        include_agenda: 是否在封面後插入目錄頁。
        include_conclusion: 是否在內容頁之後插入結論頁。
        include_closing: 是否在最後插入結尾頁。
        closing_message: 結尾頁文字。

    版面順序：封面（模板首頁）→ 目錄 →（章節分隔頁 → 內容頁…）× N
    → 結論頁 → 結尾頁。
    章節分隔頁依 ``bundle.section.chapter`` 變化自動插入；``chapter`` 為 None
    的頁面不產生分隔頁，可用來排非章節內容。

    Returns:
        :class:`RenderReport`。
    """
    template = Path(template_path or config.TEMPLATE_PPTX)

    if not template.exists():
        raise RenderError(f"找不到簡報模板：{template}")

    presentation = Presentation(str(template))
    report = RenderReport()

    if deck_title and presentation.slides:
        first_slide = presentation.slides[0]

        if first_slide.shapes.title is not None:
            text_frame = first_slide.shapes.title.text_frame
            text_frame.clear()
            run = text_frame.paragraphs[0].add_run()
            run.text = deck_title
            theme.apply_font(
                run,
                size=theme.CHAPTER_TITLE_FONT_SIZE,
                bold=True,
                color=theme.BODY_COLOR,
            )

    if not keep_template_slides:
        removed = _remove_template_placeholder_slides(presentation)

        if removed:
            report.warnings.append(
                f"已移除模板中 {removed} 張示範頁（保留首頁）"
            )

    report.chapters = chapter_order(bundles)

    # 沒有章節就不產結論頁：結論是「每個章節的結論句」收攏而成，
    # 沒有章節可收的話那一頁會是空白，比沒有這一頁更糟。
    has_conclusion = (
        include_conclusion
        and bool(report.chapters)
        and any(bundle.narrative is not None for bundle in bundles)
    )

    # 結論頁列在目錄最後一項，但不另設章節分隔頁——一張頁面的章節
    # 配一張分隔頁，讀者翻兩頁才看到一句結論。
    agenda_items = list(report.chapters)

    if has_conclusion:
        agenda_items.append(CONCLUSION_CHAPTER)

    if include_agenda and report.chapters:
        add_agenda_page(presentation, agenda_items)

    current_chapter: str | None = None
    chapter_index = 0

    for bundle in bundles:
        chapter = bundle.section.chapter

        if chapter and chapter != current_chapter:
            chapter_index += 1
            add_section_divider(presentation, chapter, index=chapter_index)
            report.divider_count += 1
            current_chapter = chapter

        try:
            _, errors = add_content_page(presentation, bundle, store)
        except RenderError as error:
            report.warnings.append(
                f"第 {bundle.section.page_number} 頁組裝失敗：{error}"
            )
            continue

        report.page_count += 1

        if bundle.chart is not None:
            report.chart_count += 1

        if errors:
            report.placeholder_errors[bundle.section.page_number or 0] = errors

    if has_conclusion:
        _, conclusion_errors = add_conclusion_page(presentation, bundles, store)
        report.conclusion_page = True

        if conclusion_errors:
            report.placeholder_errors.setdefault(0, []).extend(conclusion_errors)

    if include_closing:
        add_closing_page(presentation, closing_message)

    report.slide_count = len(presentation.slides._sldIdLst)

    target = Path(output_path or (config.OUTPUT_DIR / "deck.pptx"))
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))

    report.output_path = target
    return report


def _remove_template_placeholder_slides(presentation: Presentation) -> int:
    """
    移除模板中除首頁外的示範頁。

    python-pptx 沒有高階刪頁 API，需同時從 sldIdLst 移除項目並
    drop 對應的 relationship，否則會留下孤立的 slide part。
    """
    slide_id_list = presentation.slides._sldIdLst
    entries = list(slide_id_list)

    if len(entries) <= 1:
        return 0

    removed = 0

    for entry in entries[1:]:
        presentation.part.drop_rel(entry.rId)
        slide_id_list.remove(entry)
        removed += 1

    return removed


def scatter_labels_pending(chart: ResolvedChart) -> bool:
    """
    散點圖是否有尚未渲染的資料點標籤。

    python-pptx 無高階 API 設定散點資料點標籤文字（需操作 ``c:dLbls``），
    此函式讓上層能明確得知這項限制，而非以為標籤已經畫上去。
    見設計文件 §8.1。
    """
    return isinstance(chart.spec, ScatterSpec) and bool(chart.spec.labels)


def bundles_from_deck_spec(
    payload: dict[str, Any],
) -> tuple[list[PageBundle], MetricStore, str]:
    """Validate the renderer JSON boundary and reconstruct local objects."""
    spec = DeckSpecContract.model_validate(payload)
    store_payload = spec.metric_store.model_dump(mode="json")
    store_payload.pop("contract_version", None)
    store = MetricStore.from_dict(store_payload)
    bundles: list[PageBundle] = []

    for page in spec.pages:
        section = SectionPlan.from_dict(page.section.model_dump(mode="json"))
        plan = ChartPlan.from_dict(page.chart_plan.model_dump(mode="json"))
        chart = resolve_chart_plan(plan, store)
        narrative = PageNarrative.from_dict(
            page.narrative.model_dump(mode="json")
        )
        bundles.append(PageBundle(section, chart, narrative))

    return bundles, store, spec.title


def render_deck_from_spec(
    payload: dict[str, Any],
    *,
    output_path: str | Path,
    template_path: str | Path | None = None,
) -> RenderReport:
    """Render a deck exclusively from a schema-validated DeckSpec JSON object."""
    bundles, store, title = bundles_from_deck_spec(payload)
    return render_deck(
        bundles,
        store,
        output_path=output_path,
        template_path=template_path,
        deck_title=title,
    )

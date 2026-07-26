"""
圖表生成模組 (POC)
=====================
核心設計原則：
1. 圖表數字的「唯一真相來源」是本模組的輸入資料（解析後的 Excel 結構化資料）。
2. 一律透過 python-pptx 的 add_chart() API 生成圖表，
   讓「畫面顯示值 (chart cache)」與「右鍵編輯資料看到的內嵌工作表」
   自動保持一致，不手動操作底層 XML。
3. 不使用外部連結 (external link) 方式綁定 Excel，
   因為 PPT 寄出後外部路徑會失效。改用 PowerPoint 原生的
   「內嵌工作簿 (embedded workbook)」機制，這正是使用者
   右鍵「編輯資料」時開啟的資料來源。
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Emu


@dataclass
class ChartSpec:
    """一張圖表所需的資料規格。"""

    title: str
    categories: Sequence[str]
    series: dict[str, Sequence[float]]  # {系列名稱: 數值}
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED


def add_category_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    將 ChartSpec 的資料寫入圖表，並插入到指定投影片。

    重點：CategoryChartData 是唯一資料入口。
    python-pptx 會同時：
      1. 寫入 chart1.xml 的 <c:numCache>（畫面顯示用）
      2. 產生對應的 embedded .xlsx（右鍵「編輯資料」開啟用）
    兩者保證數值一致，因為都是從這裡的同一份 spec 產生。
    """
    chart_data = CategoryChartData()
    chart_data.categories = list(spec.categories)
    for series_name, values in spec.series.items():
        chart_data.add_series(series_name, values)

    graphic_frame = slide.shapes.add_chart(
        spec.chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = spec.title
    return chart


@dataclass
class ScatterSpec:
    """散點圖（規模 vs 成長）專用規格。"""

    title: str
    series_name: str
    points: Sequence[tuple[float, float]]  # [(x, y), ...]
    labels: Sequence[str] | None = None  # 對應每個點的標籤（如銀行名稱）


def add_scatter_chart(
    slide: Slide,
    spec: ScatterSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(8229600),
    height: int = Emu(4114800),
):
    """
    散點圖走 XyChartData，一樣是唯一資料入口原則。
    labels 部分 python-pptx 原生不支援資料點標籤文字，
    若模板要求顯示銀行名稱標籤，需要在 chart 生成後
    透過底層 XML（c:dLbls）額外處理，這裡先留 TODO。
    """
    chart_data = XyChartData()
    series = chart_data.add_series(spec.series_name)
    for x, y in spec.points:
        series.add_data_point(x, y)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = spec.title
    # TODO: 若需散點標籤（如銀行名稱），需手動操作 chart.plots[0] 的
    # dLbls XML 節點，python-pptx 目前無高階 API。
    return chart


def add_pie_chart(
    slide: Slide,
    spec: ChartSpec,
    left: int = Emu(914400),
    top: int = Emu(1600200),
    width: int = Emu(6096000),
    height: int = Emu(4114800),
):
    """
    市占率圖（圓餅圖）。走同一個 CategoryChartData 入口，
    僅 chart_type 固定為 PIE。要求 spec.series 只能有一組系列。
    """
    if len(spec.series) != 1:
        raise ValueError("圓餅圖只能有一組系列，請確認 ChartSpec.series 長度為 1")
    pie_spec = ChartSpec(
        title=spec.title,
        categories=spec.categories,
        series=spec.series,
        chart_type=XL_CHART_TYPE.PIE,
    )
    return add_category_chart(slide, pie_spec, left, top, width, height)


# ---------------------------------------------------------------------------
# Skill Registry：圖表 Agent 只透過 skill 名稱字串呼叫，不需知道實作細節。
# 新增圖表類型時，在此註冊即可，Agent 端無需改動。
# ---------------------------------------------------------------------------
CHART_SKILLS = {
    "column": add_category_chart,   # 長條圖／排名圖／成長率圖
    "bar": add_category_chart,      # 橫條圖（呼叫時 spec.chart_type 需先設為 BAR_CLUSTERED）
    "pie": add_pie_chart,           # 市占率圖
    "scatter": add_scatter_chart,   # 規模 vs 成長散點圖
    # "heatmap": 由原生表格 + 儲存格底色模擬，不走 add_chart()，
    #            另外註冊在 table_builder.py（尚未實作），Agent 需個別處理。
}

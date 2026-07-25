"""
驗證腳本 (POC)
================
目的：證明「圖表顯示值」與「右鍵編輯資料看到的內嵌工作表」數字一致。

驗證邏輯（對應規格書 T1 三方比對的簡化版）：
1. 從 chart.plots[0].categories / series 讀出「畫面顯示值」（chart XML cache）
2. 從 chart.part.chart_workbook 開啟內嵌 xlsx，讀出「編輯資料」看到的值
3. 比對兩者是否完全一致
"""

from __future__ import annotations

import io

from openpyxl import load_workbook
from pptx import Presentation

from chart_builder import ChartSpec, add_category_chart


def build_demo_pptx(template_path: str, output_path: str) -> None:
    """套用 template.pptx，在其中一頁插入一張示範圖表。"""
    prs = Presentation(template_path)

    # 對應附件三 P.5「市場規模趨勢 — 流通卡數」的示範資料
    # 這批數字模擬「Excel 解析引擎」輸出的結構化結果
    spec = ChartSpec(
        title="市場流通卡數（萬張）",
        categories=["1月", "2月", "3月", "4月", "5月", "6月",
                    "7月", "8月", "9月", "10月", "11月", "12月"],
        series={
            "市場流通卡數": [5865.5, 5874.8, 5899.4, 5938.9, 5979.6, 5986.9,
                        5994.7, 6014.9, 6035.4, 6047.9, 6042.5, 6048.6],
        },
    )

    # 使用第 4 頁（1_標題及內容 版面），對應「市場整體概況」章節下的內容頁
    slide = prs.slides[3]
    add_category_chart(slide, spec)

    prs.save(output_path)
    print(f"已生成: {output_path}")


def verify(output_path: str) -> bool:
    """比對畫面顯示值 vs 內嵌工作表值，回傳是否一致。"""
    prs = Presentation(output_path)

    all_ok = True
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            plot = chart.plots[0]
            series = plot.series[0]

            # 1) 畫面顯示值（chart XML cache）
            displayed_values = list(series.values)

            # 2) 內嵌工作表值（右鍵「編輯資料」開啟的來源）
            workbook_blob = chart.part.chart_workbook.xlsx_part.blob
            wb = load_workbook(io.BytesIO(workbook_blob), data_only=True)
            ws = wb.active
            # python-pptx 預設把資料寫在 B 欄開始（A 欄是類別）
            embedded_values = [
                cell.value for cell in ws["B"][1 : 1 + len(displayed_values)]
            ]

            match = displayed_values == embedded_values
            all_ok &= match

            print(f"[Slide {slide_idx}] 圖表: {chart.chart_title.text_frame.text}")
            print(f"  畫面顯示值 : {displayed_values}")
            print(f"  內嵌工作表值: {embedded_values}")
            print(f"  一致性: {'PASS' if match else 'FAIL'}")

    return all_ok


if __name__ == "__main__":
    TEMPLATE = "source/template.pptx"
    OUTPUT = "outpus/demo_chart.pptx"

    build_demo_pptx(TEMPLATE, OUTPUT)
    ok = verify(OUTPUT)
    print("\n=== 最終結果:", "全部一致 (PASS)" if ok else "發現不一致 (FAIL)", "===")
